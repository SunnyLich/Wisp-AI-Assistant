"""Tests for test document context contract."""

from __future__ import annotations

import importlib
import sys
from pathlib import Path
from unittest.mock import patch

import pytest


@pytest.fixture()
def win_context_fetcher():
    """Verify win context fetcher behavior."""
    with patch.object(sys, "platform", "win32"):
        import core.context_fetcher as context_fetcher

        module = importlib.reload(context_fetcher)
        try:
            yield module
        finally:
            importlib.reload(context_fetcher)


@pytest.mark.parametrize(
    ("title", "process_name", "open_path", "doc_name"),
    [
        (
            "Budget.xlsx - Excel",
            "EXCEL.EXE",
            r"C:\Users\sunny\Documents\Budget.xlsx",
            "Budget.xlsx",
        ),
        (
            "Quarterly review.pptx - PowerPoint",
            "POWERPNT.EXE",
            r"C:\Users\sunny\Documents\Quarterly review.pptx",
            "Quarterly review.pptx",
        ),
        (
            "Proposal.docx - Word Localized",
            "WINWORD.EXE",
            r"C:\Users\sunny\Documents\Proposal.docx",
            "Proposal.docx",
        ),
        (
            "Untitled 1 \u2014 LibreOffice Calc",
            "soffice.bin",
            r"C:\Users\sunny\Documents\Untitled 1.ods",
            "Untitled 1",
        ),
        (
            "Notes.odt \u2014 LibreOffice Writer",
            "soffice.bin",
            r"C:\Users\sunny\Documents\Notes.odt",
            "Notes.odt",
        ),
        (
            "laptop walmart invoice - PDF-XChange Editor",
            "PXCEditor.exe",
            r"C:\Users\sunny\Documents\laptop walmart invoice.pdf",
            "laptop walmart invoice",
        ),
        (
            "Manual.pdf - Adobe Acrobat",
            "Acrobat.exe",
            r"C:\Users\sunny\Documents\Manual.pdf",
            "Manual.pdf",
        ),
        (
            "Receipt - SumatraPDF",
            "SumatraPDF.exe",
            r"C:\Users\sunny\Documents\Receipt.pdf",
            "Receipt",
        ),
        (
            "Draft.md - Notepad++",
            "notepad++.exe",
            r"C:\Users\sunny\Documents\Draft.md",
            "Draft.md",
        ),
    ],
)
def test_common_document_apps_resolve_open_file_paths(
    win_context_fetcher,
    title,
    process_name,
    open_path,
    doc_name,
):
    """Verify common document apps resolve open file paths behavior."""
    cf = win_context_fetcher
    win = cf.WindowInfo(title=title, process_name=process_name, pid=101)

    assert cf._extract_doc_name_from_window(win) == doc_name
    with patch.object(cf, "_win_open_files_for_pid", return_value=[open_path]), \
         patch.object(cf, "_fetch_recent_files", return_value=[]):
        assert cf._resolve_doc_path(win) == open_path


def test_unsaved_common_document_window_uses_visible_text_fallback(win_context_fetcher):
    """Verify unsaved common document window uses visible text fallback behavior."""
    cf = win_context_fetcher
    cf._context_window = cf.WindowInfo(
        title="Untitled 1 \u2014 LibreOffice Calc",
        process_name="soffice.bin",
        pid=101,
        hwnd=777,
    )

    with patch.object(cf, "_enumerate_open_doc_windows", return_value=[]), \
         patch.object(cf, "_get_window_text_uia", return_value="A1\tB1\nA2\tB2"):
        assert cf.get_all_open_document_window_texts() == [
            ("Untitled 1", "A1\tB1\nA2\tB2")
        ]


def test_text_and_csv_document_readers(tmp_path):
    """Verify text and csv document readers behavior."""
    from core.llm_clients import client as llm

    text_path = tmp_path / "notes.md"
    text_path.write_text("# Notes\nhello text", encoding="utf-8")
    csv_path = tmp_path / "table.csv"
    csv_path.write_text("name,value\nalpha,1", encoding="utf-8")

    assert "hello text" in llm.read_document_file(str(text_path))
    assert "alpha,1" in llm.read_document_file(str(csv_path))


def test_docx_reader_extracts_paragraph_text(tmp_path):
    """Verify docx reader extracts paragraph text behavior."""
    pytest.importorskip("docx")
    from docx import Document
    from core.llm_clients import client as llm

    path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_paragraph("Docx contract text")
    doc.save(path)

    assert "Docx contract text" in llm.read_document_file(str(path))


def test_xlsx_reader_extracts_sheet_cells(tmp_path):
    """Verify xlsx reader extracts sheet cells behavior."""
    pytest.importorskip("openpyxl")
    import openpyxl
    from core.llm_clients import client as llm

    path = tmp_path / "sample.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Invoices"
    ws["A1"] = "Item"
    ws["B1"] = "Total"
    ws["A2"] = "Laptop"
    ws["B2"] = 999
    wb.save(path)

    text = llm.read_document_file(str(path))

    assert "[Sheet: Invoices]" in text
    assert "Laptop\t999" in text


def test_pptx_reader_extracts_slide_text(tmp_path):
    """Verify pptx reader extracts slide text behavior."""
    pytest.importorskip("pptx")
    from pptx import Presentation
    from core.llm_clients import client as llm

    path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Deck contract title"
    prs.save(path)

    assert "Deck contract title" in llm.read_document_file(str(path))


def test_odt_reader_extracts_paragraph_text(tmp_path):
    """Verify odt reader extracts paragraph text behavior."""
    pytest.importorskip("odf")
    from odf import text as odf_text
    from odf.opendocument import OpenDocumentText
    from core.llm_clients import client as llm

    path = tmp_path / "sample.odt"
    doc = OpenDocumentText()
    doc.text.addElement(odf_text.P(text="ODT contract text"))
    doc.save(str(path))

    assert "ODT contract text" in llm.read_document_file(str(path))


def test_pdf_reader_dispatches_to_pdf_text_extractor(tmp_path, monkeypatch):
    """Verify pdf reader dispatches to pdf text extractor behavior."""
    from core.llm_clients import client as llm
    from core.llm_clients import documents

    path = tmp_path / "sample.pdf"
    path.write_bytes(b"%PDF-1.4\n% test placeholder\n")
    # PDF extraction lives in core.llm_clients.documents (re-exported via client);
    # patch it at the seam where read_document_file resolves it.
    monkeypatch.setattr(documents, "_read_pdf_text", lambda p, _max_chars: f"PDF text from {Path(p).name}")

    assert "PDF text from sample.pdf" in llm.read_document_file(str(path))
