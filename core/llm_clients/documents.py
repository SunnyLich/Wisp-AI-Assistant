"""Local document reading for LLM context.

Extracts plain text from documents (PDF, Office, OpenDocument, plain-text)
for injection into prompts and the read_document tool. Split out of
core.llm_clients.client; the names are re-exported there for compatibility.
"""
from __future__ import annotations

import os
import zipfile

import config
from core.llm_clients.logging_utils import log_context as _log_context

_MAX_DOCUMENT_INPUT_BYTES = 25 * 1024 * 1024
_MAX_ARCHIVE_EXPANDED_BYTES = 100 * 1024 * 1024
_MAX_ARCHIVE_ENTRIES = 10_000
_ARCHIVE_DOCUMENT_EXTS = {".docx", ".xlsx", ".pptx", ".odt", ".ods", ".odp"}


def _document_safety_error(path: str, ext: str) -> str:
    """Return a small error for sources unsafe to expand, otherwise empty."""
    try:
        size = os.path.getsize(path)
    except OSError as exc:
        return f"Failed to read {path!r}: {exc}"
    if size > _MAX_DOCUMENT_INPUT_BYTES:
        return (
            f"Failed to read {path!r}: file is {size:,} bytes; "
            f"the capture safety limit is {_MAX_DOCUMENT_INPUT_BYTES:,} bytes."
        )
    if ext not in _ARCHIVE_DOCUMENT_EXTS or not zipfile.is_zipfile(path):
        return ""
    try:
        with zipfile.ZipFile(path) as archive:
            infos = archive.infolist()
            expanded = sum(max(0, int(info.file_size)) for info in infos)
    except (OSError, zipfile.BadZipFile) as exc:
        return f"Failed to read {path!r}: invalid document archive ({exc})"
    if len(infos) > _MAX_ARCHIVE_ENTRIES or expanded > _MAX_ARCHIVE_EXPANDED_BYTES:
        return (
            f"Failed to read {path!r}: expanded document is too large "
            f"({expanded:,} bytes across {len(infos):,} entries)."
        )
    return ""


def _ambient_document_max_chars() -> int:
    """Handle ambient document max chars for LLM clients client."""
    return config.get_settings().context.ambient_document_max_chars

def _tool_document_max_chars() -> int:
    """Handle tool document max chars for LLM clients client."""
    return config.get_settings().context.tool_document_max_chars

def _normalize_pdf_text(s: str) -> str:
    """Collapse layout whitespace from PDF text.

    LiteParse pads its output with horizontal spacing used for visual
    layout, which carries no extra content but multiplies the token count
    sent to the LLM. Collapsing intra-line whitespace yields the same text
    pypdf would, at a fraction of the tokens.
    """
    import re
    lines = []
    for line in s.splitlines():
        line = re.sub(r"[ \t]+", " ", line).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)

def _read_pdf_text(path: str, max_chars: int) -> str:
    """Extract PDF text, preferring LiteParse (fast, native) over pypdf."""
    parts: list[str] = []
    total = 0
    try:
        import liteparse  # type: ignore
    except ImportError:
        liteparse = None
    if liteparse is not None:
        try:
            # LiteParse parses every page up front, so cap pages to roughly
            # what max_chars can hold (with a buffer for sparse pages) instead
            # of parsing the whole document just to truncate the output.
            page_cap = max(8, max_chars // 500 + 5)
            lp = liteparse.LiteParse(ocr_enabled=False, quiet=True, max_pages=page_cap)
            result = lp.parse(path)
            for i in range(1, result.num_pages + 1):
                page = result.get_page(i)
                page_text = _normalize_pdf_text(page.text) if page and page.text else ""
                if page_text:
                    parts.append(f"[Page {i}]\n{page_text}")
                    total += len(page_text)
                    if total > max_chars:
                        break
            return "\n\n".join(parts)
        except Exception:
            parts.clear()
    # Fallback: pure-Python pypdf (slower, no native dependency).
    import pypdf  # type: ignore
    reader = pypdf.PdfReader(path)
    for i, page in enumerate(reader.pages, 1):
        page_text = (page.extract_text() or "").strip()
        if page_text:
            parts.append(f"[Page {i}]\n{page_text}")
            total += len(page_text)
            if total > max_chars:
                break
    return "\n\n".join(parts)

def _read_document_file(path: str, max_chars: int | None = None) -> str:
    """Read a local document file and return its plain text."""
    if max_chars is None:
        max_chars = _ambient_document_max_chars()
    ext = os.path.splitext(path)[1].lower()
    try:
        safety_error = _document_safety_error(path, ext)
        if safety_error:
            return safety_error
        if ext == ".docx":
            from docx import Document  # type: ignore
            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext in (".xlsx", ".xls"):
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        parts.append("\t".join(cells))
            text = "\n".join(parts)
        elif ext == ".pptx":
            from pptx import Presentation  # type: ignore
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                parts.append(line)
            text = "\n".join(parts)
        elif ext == ".pdf":
            text = _read_pdf_text(path, max_chars)
        elif ext in (".odt", ".ods", ".odp"):
            from odf import teletype
            from odf import text as odf_text  # type: ignore
            from odf.opendocument import load as odf_load  # type: ignore
            doc = odf_load(path)
            paragraphs = doc.getElementsByType(odf_text.P)
            text = "\n".join(
                teletype.extractText(p) for p in paragraphs
                if teletype.extractText(p).strip()
            )
        elif ext in (".txt", ".md", ".csv", ".py", ".js", ".ts",
                     ".json", ".xml", ".html", ".log"):
            with open(path, encoding="utf-8", errors="replace") as f:
                # Never materialize a giant text file just to discard its tail.
                text = f.read(max(1, max_chars) + 1)
        else:
            return f"File type {ext!r} is not supported for reading."
        if len(text) > max_chars:
            text = text[:max_chars] + "\n[-¦truncated]"
        # Redact sensitive data before the text reaches the LLM.
        from core.context_fetcher import _redact  # noqa: PLC0415
        text = _redact(text)
        _log_context(f"tool: read_document - read {path!r}", text)
        return text
    except Exception as e:
        return f"Failed to read {path!r}: {e}"

def _read_document_paths(
    paths: list[str],
    max_chars_per_doc: int | None = None,
) -> str:
    """Read multiple local document files and join readable results."""
    import os

    if max_chars_per_doc is None:
        max_chars_per_doc = _ambient_document_max_chars()
    parts: list[str] = []
    for path in paths:
        text = _read_document_file(path, max_chars=max_chars_per_doc)
        if text and not text.startswith(("Could not", "File type", "Failed to")):
            parts.append(f"[{os.path.basename(path)}]\n{text}")
    return "\n\n".join(parts)
