"""
core/llm.py -” Cloud LLM client with streaming support.

Supports:
  - Groq (OpenAI-compatible, fast TTFT)
  - OpenAI
  - Anthropic Claude

Clients are module-level singletons so the TLS connection is reused across
calls, eliminating handshake overhead from every request.
"""
from __future__ import annotations
import gzip
import json as _stdlib_json
import ssl as _ssl
import threading as _threading
import urllib.error as _urllib_error
import urllib.request as _urllib_request
import config
from pathlib import Path
from core.tool_registry import ToolRegistry, ToolSpec
from core.tools.local_files import execute_live_file_tool, normalize_file_access_mode
from core.system import macos_safety
from core.system.native_locks import native_init_lock, ssl_init_lock
from core.system import sdk_clients
from core.llm_clients.routes import (
    GOOGLE_OPENAI_BASE_URL as _GOOGLE_OPENAI_BASE_URL,
    DEEPSEEK_BASE_URL as _DEEPSEEK_BASE_URL,
    OPENROUTER_BASE_URL as _OPENROUTER_BASE_URL,
    MISTRAL_BASE_URL as _MISTRAL_BASE_URL,
    XAI_BASE_URL as _XAI_BASE_URL,
    TOGETHER_BASE_URL as _TOGETHER_BASE_URL,
    CEREBRAS_BASE_URL as _CEREBRAS_BASE_URL,
    OLLAMA_BASE_URL as _OLLAMA_BASE_URL,
    api_key_for as _api_key_for,
    credential_source_for_provider as _credential_source_for_provider,
    parse_model_fallbacks as _parse_model_fallbacks,
    route_candidates as _route_candidates,
    normalize_model_for_provider as _normalize_model_for_provider,
)
from core.llm_clients.logging_utils import log_event
from core.llm_clients.messages import (
    build_openai_messages as _build_openai_messages,
    sanitize_history as _sanitize_history,
)
from core.llm_clients.prompt_guidance import (
    REWRITE_SYSTEM_PROMPT as _REWRITE_SYSTEM_PROMPT,
    with_memory_save_note as _with_memory_save_note,
    with_memory_search_note as _with_memory_search_note,
    with_screenshot_note as _with_screenshot_note,
    with_tools_note as _with_tools_note,
)
from dataclasses import dataclass, field
from typing import Callable, Generator

_TOOL_REGISTRY = ToolRegistry()
_LOCAL_FILE_TOOLS = {"list_files", "read_file", "create_file", "edit_file", "write_file"}
_FILE_EDIT_APPROVAL_CALLBACK: Callable[[dict], bool] | None = None
_LIVE_TOOL_CONTEXT = _threading.local()


def _log_context(
    reason: str,
    text: str,
    max_line: int = 120,
    max_lines: int = 12,
    max_chars: int = 1200,
) -> None:
    """Log a compact preview of a context block for debugging."""

    def _trim(line: str) -> str:
        """Handle trim for LLM clients client."""
        return line if len(line) <= max_line else line[:max_line] + "-¦"

    lines = [_trim(l) for l in text.splitlines() if l.strip()]
    truncated = False

    if len(lines) > max_lines:
        lines = lines[:max_lines]
        truncated = True

    body = "\n  ".join(lines) if lines else "[empty]"
    if len(body) > max_chars:
        body = body[:max_chars].rstrip() + "-¦"
        truncated = True
    if truncated and body != "[empty]":
        body += "\n  [preview truncated]"

    log_event(
        "llm.context_preview",
        f"Context preview for {reason}:\n  {body}",
        reason=reason,
        preview=body,
        truncated=truncated,
    )


def _ambient_document_max_chars() -> int:
    """Handle ambient document max chars for LLM clients client."""
    return config.get_settings().context.ambient_document_max_chars


def _tool_document_max_chars() -> int:
    """Handle tool document max chars for LLM clients client."""
    return config.get_settings().context.tool_document_max_chars


def _normalize_pdf_text(s: str) -> str:
    """Collapse layout whitespace from PDF text.

    LiteParse pads its output with horizontal spacing used for visual
    layout, which carries no extra content but multiplies the token count
    sent to the LLM. Collapsing intra-line whitespace yields the same text
    pypdf would, at a fraction of the tokens.
    """
    import re
    lines = []
    for line in s.splitlines():
        line = re.sub(r"[ \t]+", " ", line).strip()
        if line:
            lines.append(line)
    return "\n".join(lines)


def _read_pdf_text(path: str, max_chars: int) -> str:
    """Extract PDF text, preferring LiteParse (fast, native) over pypdf."""
    parts: list[str] = []
    total = 0
    try:
        import liteparse  # type: ignore
    except ImportError:
        liteparse = None
    if liteparse is not None:
        try:
            # LiteParse parses every page up front, so cap pages to roughly
            # what max_chars can hold (with a buffer for sparse pages) instead
            # of parsing the whole document just to truncate the output.
            page_cap = max(8, max_chars // 500 + 5)
            lp = liteparse.LiteParse(ocr_enabled=False, quiet=True, max_pages=page_cap)
            result = lp.parse(path)
            for i in range(1, result.num_pages + 1):
                page = result.get_page(i)
                page_text = _normalize_pdf_text(page.text) if page and page.text else ""
                if page_text:
                    parts.append(f"[Page {i}]\n{page_text}")
                    total += len(page_text)
                    if total > max_chars:
                        break
            return "\n\n".join(parts)
        except Exception:
            parts.clear()
    # Fallback: pure-Python pypdf (slower, no native dependency).
    import pypdf  # type: ignore
    reader = pypdf.PdfReader(path)
    for i, page in enumerate(reader.pages, 1):
        page_text = (page.extract_text() or "").strip()
        if page_text:
            parts.append(f"[Page {i}]\n{page_text}")
            total += len(page_text)
            if total > max_chars:
                break
    return "\n\n".join(parts)


def _read_document_file(path: str, max_chars: int | None = None) -> str:
    """Read a local document file and return its plain text."""
    import os
    if max_chars is None:
        max_chars = _ambient_document_max_chars()
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".docx":
            from docx import Document  # type: ignore
            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext in (".xlsx", ".xls"):
            import openpyxl  # type: ignore
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                parts.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        parts.append("\t".join(cells))
            text = "\n".join(parts)
        elif ext == ".pptx":
            from pptx import Presentation  # type: ignore
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                parts.append(line)
            text = "\n".join(parts)
        elif ext == ".pdf":
            text = _read_pdf_text(path, max_chars)
        elif ext in (".odt", ".ods", ".odp"):
            from odf import text as odf_text, teletype  # type: ignore
            from odf.opendocument import load as odf_load  # type: ignore
            doc = odf_load(path)
            paragraphs = doc.getElementsByType(odf_text.P)
            text = "\n".join(
                teletype.extractText(p) for p in paragraphs
                if teletype.extractText(p).strip()
            )
        elif ext in (".txt", ".md", ".csv", ".py", ".js", ".ts",
                     ".json", ".xml", ".html", ".log"):
            with open(path, encoding="utf-8", errors="replace") as f:
                text = f.read()
        else:
            return f"File type {ext!r} is not supported for reading."
        if len(text) > max_chars:
            text = text[:max_chars] + "\n[-¦truncated]"
        # Redact sensitive data before the text reaches the LLM.
        from core.context_fetcher import _redact  # noqa: PLC0415
        text = _redact(text)
        _log_context(f"tool: read_document -” read {path!r}", text)
        return text
    except Exception as e:
        return f"Failed to read {path!r}: {e}"


def _read_document_paths(
    paths: list[str],
    max_chars_per_doc: int | None = None,
) -> str:
    """Read multiple local document files and join readable results."""
    import os

    if max_chars_per_doc is None:
        max_chars_per_doc = _ambient_document_max_chars()
    parts: list[str] = []
    for path in paths:
        text = _read_document_file(path, max_chars=max_chars_per_doc)
        if text and not text.startswith(("Could not", "File type", "Failed to")):
            parts.append(f"[{os.path.basename(path)}]\n{text}")
    return "\n\n".join(parts)


def _execute_get_context(inputs: dict) -> str:
    """Built-in context tool: fetch a browser page or open document text."""
    url = inputs.get("url", "").strip()
    if url:
        from core.context_fetcher import fetch_browser_content_for_tool
        result = fetch_browser_content_for_tool(url)
        _log_context(
            f"tool: get_context (browser) -” {url!r}",
            result or "",
        )
        return result or f"Could not fetch content from {url!r}."

    from core.context_fetcher import get_all_open_document_paths

    paths = get_all_open_document_paths()
    if not paths:
        return "Could not determine any open document paths from supported app windows."

    text = _read_document_paths(paths, max_chars_per_doc=_tool_document_max_chars())
    if text:
        return text
    return "Open document windows were detected, but none of their file types were readable."


def _execute_memory_search(inputs: dict) -> str:
    """Handle execute memory search for LLM clients client."""
    query = str((inputs or {}).get("query") or "").strip()
    top_k_raw = (inputs or {}).get("top_k")
    try:
        top_k = int(top_k_raw) if top_k_raw is not None else None
    except (TypeError, ValueError):
        top_k = None
    if not query:
        return "Memory search requires a query."
    try:
        from core.memory_store import store

        return store.get_manager().retrieve_relevant(query, top_k=top_k) or "No relevant memory found."
    except Exception as exc:  # noqa: BLE001 - memory should not block the answer path
        return f"Memory search failed: {type(exc).__name__}: {exc}"


def _execute_memory_save(inputs: dict) -> str:
    """Handle execute memory save for LLM clients client."""
    text = str((inputs or {}).get("text") or "").strip()
    scope = str((inputs or {}).get("scope") or "").strip().lower()
    # Unset/unknown scope means "follow the conversation's project" (the store
    # default); only an explicit "general" promotes a fact to global.
    if scope not in ("general", "project"):
        scope = ""
    if not text:
        return "Memory save requires the fact text to store."
    try:
        from core.memory_store import store

        result = store.get_manager().save_memory(text, scope=scope or None)
        if not result.get("ok"):
            return f"Did not store memory: {result.get('reason', 'rejected')}."
        where = "the current project" if result.get("scope") == "project" else "general memory"
        return f"Stored to {where}: {result.get('text', text)!r}"
    except Exception as exc:  # noqa: BLE001 - memory should not block the answer path
        return f"Memory save failed: {type(exc).__name__}: {exc}"


def _execute_git_status(inputs: dict) -> str:
    """Handle execute git status for LLM clients client."""
    cwd = inputs.get("cwd") or config.TOOL_GIT_ROOT
    return _run_read_only_command(["git", "status", "--short"], cwd=cwd)


def _execute_git_diff(inputs: dict) -> str:
    """Handle execute git diff for LLM clients client."""
    cwd = inputs.get("cwd") or config.TOOL_GIT_ROOT
    return _run_read_only_command(["git", "diff", "--", "."], cwd=cwd)


def _run_read_only_command(args: list[str], cwd: str) -> str:
    """Run read only command."""
    import subprocess
    from pathlib import Path

    root = Path(cwd or ".").expanduser().resolve()
    completed = subprocess.run(
        args,
        cwd=str(root),
        capture_output=True,
        text=True,
        timeout=20,
    )
    output = (completed.stdout or completed.stderr or "").strip()
    return output[:12000] or f"{' '.join(args)} returned no output."


def set_file_edit_approval_callback(callback: Callable[[dict], bool] | None) -> None:
    """Register the UI callback used by ask-mode local file edits."""
    global _FILE_EDIT_APPROVAL_CALLBACK
    _FILE_EDIT_APPROVAL_CALLBACK = callback


def set_live_file_access_mode(mode: str | None) -> None:
    """Set local-file access mode for the current live model request thread."""
    if mode is None:
        if hasattr(_LIVE_TOOL_CONTEXT, "file_access_mode"):
            delattr(_LIVE_TOOL_CONTEXT, "file_access_mode")
        return
    _LIVE_TOOL_CONTEXT.file_access_mode = normalize_file_access_mode(mode)


def set_live_file_approval_callback(callback: Callable[[dict], bool] | None) -> None:
    """Set the approval callback for the current live model request thread."""
    if callback is None:
        if hasattr(_LIVE_TOOL_CONTEXT, "file_approval_callback"):
            delattr(_LIVE_TOOL_CONTEXT, "file_approval_callback")
        return
    _LIVE_TOOL_CONTEXT.file_approval_callback = callback


def _effective_live_file_approval_callback() -> Callable[[dict], bool] | None:
    """Return the live request approval callback, falling back to the global hook."""
    return getattr(_LIVE_TOOL_CONTEXT, "file_approval_callback", None) or _FILE_EDIT_APPROVAL_CALLBACK


def _execute_list_files(inputs: dict) -> str:
    """List text-accessible files within a configured local-file root."""
    return execute_live_file_tool(
        "list_files",
        inputs or {},
        access_mode=_effective_live_file_access_mode(["list_files"]),
        approval_callback=_effective_live_file_approval_callback(),
    )


def _execute_read_file(inputs: dict) -> str:
    """Read a UTF-8 text file within a configured local-file root."""
    return execute_live_file_tool(
        "read_file",
        inputs or {},
        access_mode=_effective_live_file_access_mode(["read_file"]),
        approval_callback=_effective_live_file_approval_callback(),
    )


def _execute_edit_file(inputs: dict) -> str:
    """Replace one exact text span in a configured local-file root."""
    return execute_live_file_tool(
        "edit_file",
        inputs or {},
        access_mode=_effective_live_file_access_mode(["edit_file"]),
        approval_callback=_effective_live_file_approval_callback(),
    )


def _execute_create_file(inputs: dict) -> str:
    """Create a new text file within a configured local-file root."""
    return execute_live_file_tool(
        "create_file",
        inputs or {},
        access_mode=_effective_live_file_access_mode(["create_file"]),
        approval_callback=_effective_live_file_approval_callback(),
    )


def _execute_write_file(inputs: dict) -> str:
    """Create or overwrite a text file in a configured local-file root."""
    return execute_live_file_tool(
        "write_file",
        inputs or {},
        access_mode=_effective_live_file_access_mode(["write_file"]),
        approval_callback=_effective_live_file_approval_callback(),
    )


def _execute_github_repo(inputs: dict) -> str:
    """Handle execute github repo for LLM clients client."""
    repo = str(inputs.get("repo") or "").strip()
    if not repo:
        return "Missing repo. Use owner/name."
    return _github_get_json(f"https://api.github.com/repos/{repo}")


def _execute_github_issue(inputs: dict) -> str:
    """Handle execute github issue for LLM clients client."""
    repo = str(inputs.get("repo") or "").strip()
    number = str(inputs.get("number") or "").strip()
    if not repo or not number:
        return "Missing repo or number."
    return _github_get_json(f"https://api.github.com/repos/{repo}/issues/{number}")


def _github_get_json(url: str) -> str:
    """Handle github get json for LLM clients client."""
    import json
    import urllib.request
    from core.auth import github as github_auth

    token = github_auth.get_valid_access_token()
    if not token:
        return "GitHub OAuth is not configured. Sign in from Settings."
    req = urllib.request.Request(
        url,
        headers={
            "Accept": "application/vnd.github+json",
            "Authorization": f"Bearer {token}",
            "User-Agent": "python-ai-overlay",
            "X-GitHub-Api-Version": "2022-11-28",
        },
    )
    with urllib.request.urlopen(req, timeout=20) as resp:
        data = json.loads(resp.read())
    return json.dumps(data, indent=2, ensure_ascii=False)[:12000]


def _register_builtin_tools() -> None:
    """Handle register builtin tools for LLM clients client."""
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="web_search",
            description="Search the web for current information.",
            input_schema={"type": "object", "properties": {}, "required": []},
            server_schema={
                "type": "web_search_20250305",
                "name": "web_search",
                "max_uses": 2,
            },
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="get_context",
            description=(
                "Retrieve additional context the user can see. "
                "Pass a URL to fetch a web page; omit it to read open local "
                "documents from supported apps "
                "(Word, Excel, PowerPoint, PDF, LibreOffice, Notepad, etc.)."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "url": {
                        "type": "string",
                        "description": (
                            "A web page URL (http:// or https://) to fetch. "
                            "Omit this field to read open local documents instead."
                        ),
                    }
                },
                "required": [],
            },
            executor=_execute_get_context,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="capture_screen",
            description=(
                "Take a screenshot of the user's primary monitor and see it. "
                "Use this only when the text context already provided is not "
                "enough to answer and you need to visually see what is on the "
                "user's screen (UI layout, a chart, an image, an error dialog, "
                "etc.)."
            ),
            input_schema={"type": "object", "properties": {}, "required": []},
            # Handled directly in the Anthropic and OpenAI tool loops, which
            # deliver the screenshot as image content. The executor below is a
            # text fallback for any path that calls it without that handling.
            executor=lambda _inputs: (
                "Screen capture could not be returned in this context."
            ),
            opt_in=True,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="memory_search",
            description=(
                "Search the user's stored long-term memory for relevant facts. "
                "Use this only when personal or project memory would help answer."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "Short search query for the memory facts to retrieve.",
                    },
                    "top_k": {
                        "type": "integer",
                        "description": "Optional maximum number of facts to return.",
                    },
                },
                "required": ["query"],
            },
            executor=_execute_memory_search,
            opt_in=True,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="memory_save",
            description=(
                "Save a durable fact about the user to long-term memory so it is "
                "available in future conversations. Use this whenever the user "
                "shares a stable preference, personal detail, or project fact "
                "worth remembering (not transient or one-off requests). By "
                "default a saved fact is scoped to the conversation's current "
                "project; set scope='general' only for universal facts (like a "
                "personal preference) that should apply across every project."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "text": {
                        "type": "string",
                        "description": (
                            "The single durable fact to remember, as a short "
                            "self-contained sentence."
                        ),
                    },
                    "scope": {
                        "type": "string",
                        "enum": ["general", "project"],
                        "description": (
                            "Omit to scope the fact to the current project "
                            "(the default). Use 'general' to promote a universal "
                            "fact that applies everywhere; 'project' is the "
                            "explicit form of the default."
                        ),
                    },
                },
                "required": ["text"],
            },
            executor=_execute_memory_save,
            opt_in=True,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="git_status",
            description="Return read-only git status for the configured local repository.",
            input_schema={
                "type": "object",
                "properties": {"cwd": {"type": "string", "description": "Optional local repo path."}},
                "required": [],
            },
            executor=_execute_git_status,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="git_diff",
            description="Return read-only git diff for the configured local repository.",
            input_schema={
                "type": "object",
                "properties": {"cwd": {"type": "string", "description": "Optional local repo path."}},
                "required": [],
            },
            executor=_execute_git_diff,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="list_files",
            description=(
                "List files under a folder inside the user's configured local-file roots. "
                "Requires the caller to explicitly allow local file tools."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "root": {"type": "string", "description": "Required for relative folders when multiple roots are configured."},
                    "folder": {"type": "string", "description": "Folder to list. Use an absolute path or a path relative to root."},
                    "limit": {"type": "integer", "description": "Maximum files to return."},
                },
                "required": [],
            },
            executor=_execute_list_files,
            opt_in=True,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="read_file",
            description=(
                "Read a UTF-8 text file inside the user's configured local-file roots. "
                "Blocked globs and size limits are enforced."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "root": {"type": "string", "description": "Required for relative paths when multiple roots are configured."},
                    "path": {"type": "string", "description": "Absolute path, or a path relative to root."},
                    "max_chars": {"type": "integer", "description": "Maximum characters to return."},
                },
                "required": ["path"],
            },
            executor=_execute_read_file,
            opt_in=True,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="edit_file",
            description=(
                "Replace exactly one matching text span in a file inside the user's "
                "configured local-file roots. Uses the configured never/ask/auto edit mode."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "root": {"type": "string", "description": "Required for relative paths when multiple roots are configured."},
                    "path": {"type": "string", "description": "Absolute path, or a path relative to root."},
                    "old": {"type": "string", "description": "Exact existing text to replace. Must match exactly once."},
                    "new": {"type": "string", "description": "Replacement text."},
                },
                "required": ["path", "old", "new"],
            },
            executor=_execute_edit_file,
            opt_in=True,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="create_file",
            description=(
                "Create a new UTF-8 text file inside the user's configured "
                "local-file roots. Fails if the file already exists. Uses the "
                "configured never/ask/auto create mode."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "root": {"type": "string", "description": "Required for relative paths when multiple roots are configured."},
                    "path": {"type": "string", "description": "Absolute path, or a path relative to root."},
                    "content": {"type": "string", "description": "Full file content to write."},
                },
                "required": ["path", "content"],
            },
            executor=_execute_create_file,
            opt_in=True,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="write_file",
            description=(
                "Create or overwrite a UTF-8 text file inside the user's configured "
                "local-file roots. Uses the configured never/ask/auto edit mode."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "root": {"type": "string", "description": "Required for relative paths when multiple roots are configured."},
                    "path": {"type": "string", "description": "Absolute path, or a path relative to root."},
                    "content": {"type": "string", "description": "Full file content to write."},
                },
                "required": ["path", "content"],
            },
            executor=_execute_write_file,
            opt_in=True,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="github_repo",
            description="Fetch GitHub repository metadata using the signed-in GitHub OAuth account.",
            input_schema={
                "type": "object",
                "properties": {"repo": {"type": "string", "description": "Repository in owner/name form."}},
                "required": ["repo"],
            },
            executor=_execute_github_repo,
        )
    )
    _TOOL_REGISTRY.register_builtin(
        ToolSpec(
            name="github_issue",
            description="Fetch a GitHub issue or pull request by repository and number.",
            input_schema={
                "type": "object",
                "properties": {
                    "repo": {"type": "string", "description": "Repository in owner/name form."},
                    "number": {"type": "integer", "description": "Issue or pull request number."},
                },
                "required": ["repo", "number"],
            },
            executor=_execute_github_issue,
        )
    )


def _looks_like_screenshot_tool_request(text: str) -> bool:
    """Detect providers that ask for the screenshot tool as text, not tool_calls."""
    normalized = " ".join((text or "").lower().split())
    if not normalized:
        return False
    if "capture_screen" in normalized or "screenshot tool" in normalized:
        return True
    visual_word = any(
        word in normalized
        for word in ("screenshot", "screen", "visual", "image", "see what")
    )
    request_word = any(
        word in normalized
        for word in (
            "need",
            "needs",
            "use",
            "take",
            "capture",
            "access",
            "can't see",
            "cannot see",
            "only see",
        )
    )
    return visual_word and request_word


# Substrings of model names known to accept image input. Best-effort only, used
# for a settings-time heads-up; unknown models are treated as text-only so we err
# toward warning. This never gates runtime behavior.
_VISION_MODEL_HINTS = (
    "gpt-4o", "gpt-4.1", "gpt-4-turbo", "gpt-4-vision", "gpt-5", "chatgpt-4o",
    "o1", "o3", "o4",
    "claude", "sonnet", "opus", "haiku",
    "gemini",
    "pixtral", "mistral-small-3", "mistral-medium",
    "llama-3.2-11b", "llama-3.2-90b", "llama-4", "scout", "maverick",
    "qwen-vl", "qwen2-vl", "qwen2.5-vl", "qwen3-vl",
    "internvl", "phi-3.5-vision", "phi-4-multimodal",
    "grok-2-vision", "grok-4",
    "vision", "-vl", "multimodal",
)


def _model_accepts_images(model: str) -> bool:
    """Handle model accepts images for LLM clients client."""
    m = (model or "").strip().lower()
    return any(hint in m for hint in _VISION_MODEL_HINTS)


# Model substrings whose providers reject any non-default sampling value
# (temperature / top_p): GPT-5 family + OpenAI o-series reasoning models, and the
# newest Claude models (Opus 4.7+, Fable), which removed those parameters. For
# these we omit temperature and let the model use its default.
_NO_CUSTOM_SAMPLING_HINTS = (
    "gpt-5", "o1", "o3", "o4",
    "opus-4-7", "opus-4-8", "fable",
)


def _model_rejects_custom_sampling(model: str) -> bool:
    """Handle model rejects custom sampling for LLM clients client."""
    m = (model or "").strip().lower()
    return any(hint in m for hint in _NO_CUSTOM_SAMPLING_HINTS)


def _apply_sampling(kwargs: dict, model: str, temperature: float | None) -> dict:
    """Add ``temperature`` only when the model accepts a custom one.

    Comply with each model's rules up front: GPT-5/o-series and the newest Claude
    models only accept their default sampling value and 400 on anything else, so
    we omit it for them. Models that do accept it keep the requested value. The
    reactive drop-and-retry on the OpenAI-compatible route remains the backstop
    for any model not covered by the hint list above.
    """
    if temperature is not None and not _model_rejects_custom_sampling(model):
        kwargs["temperature"] = temperature
    return kwargs


# OpenAI's GPT-5 family and o-series reasoning models reject ``max_tokens`` and
# require ``max_completion_tokens``. Only OpenAI serves these model names, so the
# substring match is effectively provider-scoped.
_MAX_COMPLETION_TOKENS_HINTS = ("gpt-5", "o1", "o3", "o4")


def _model_uses_max_completion_tokens(model: str) -> bool:
    """Handle model uses max completion tokens for LLM clients client."""
    m = (model or "").strip().lower()
    return any(hint in m for hint in _MAX_COMPLETION_TOKENS_HINTS)


def _apply_max_output(kwargs: dict, model: str, value) -> dict:
    """Set the output-token cap under the field name the model accepts.

    GPT-5 / o-series want ``max_completion_tokens``; everything else takes
    ``max_tokens``. Complying up front avoids a 400 + retry round-trip.
    """
    if value is None:
        return kwargs
    key = "max_completion_tokens" if _model_uses_max_completion_tokens(model) else "max_tokens"
    kwargs[key] = value
    return kwargs


def screenshot_capability_warnings(
    screenshot_modes,
    *,
    llm_provider: str,
    llm_model: str,
    vision_provider: str,
    vision_model: str,
) -> list[str]:
    """Best-effort warnings if enabled screenshot modes likely can't be served.

    Advisory only — the caller should still honor the user's setting and just
    surface these. *screenshot_modes* is the per-caller list of "off"/"auto"/
    "model" values. Returns human-readable warning strings (possibly empty).
    """
    modes = {m for m in screenshot_modes if m and m != "off"}
    if not modes:
        return []

    llm_provider = (llm_provider or "").strip().lower()
    vision_provider = (vision_provider or "").strip().lower()
    vision_model = (vision_model or "").strip()
    warnings: list[str] = []

    if "auto" in modes:
        # Eager ("On") screenshots always route through the Vision LLM.
        if not vision_model:
            warnings.append(
                "Auto screenshot needs an Image model. Pick one under Image model."
            )
        elif vision_provider == "copilot":
            warnings.append(
                "Copilot cannot read screenshots. Pick a different Image model."
            )
        elif not _model_accepts_images(vision_model):
            warnings.append(
                "This Image model may not read screenshots. Pick a model with image support."
            )

    if "model" in modes:
        # On-demand ("Let model decide") screenshots run through the tool loop.
        if llm_provider in ("copilot", "chatgpt"):
            warnings.append(
                "'Let model decide' screenshots do not work with ChatGPT or Copilot. "
                "Pick a different Chat model, or turn screenshots On instead."
            )
        elif llm_provider == "anthropic":
            pass  # Claude tool models accept images
        elif llm_provider in _OPENAI_COMPAT_PROVIDER_SET:
            # Answered by a same-provider Vision model if set, else the main model.
            target = (
                vision_model
                if (vision_provider == llm_provider and vision_model)
                else (llm_model or "").strip()
            )
            if not _model_accepts_images(target):
                warnings.append(
                    "This Chat model may not read screenshots. Pick an image-capable "
                    "Chat model, or use Image model with the same provider."
                )

    return warnings


def tool_capability_warnings(tools_enabled: bool, *, llm_provider: str) -> list[str]:
    """Best-effort warning if the context-tools setting won't behave as a user
    expects on the active provider. Advisory only."""
    if not tools_enabled:
        return []
    if (llm_provider or "").strip().lower() == "chatgpt":
        return [
            "ChatGPT cannot use live context tools here. Use On to attach context "
            "up front, or choose another Chat model for Let model decide."
        ]
    return []


# Providers that authenticate with a personal subscription login (OAuth) instead
# of a static API key. Their tokens can be invalidated server-side and need
# periodic re-sign-in; API keys persist indefinitely.
_SUBSCRIPTION_AUTH_PROVIDERS = frozenset({"chatgpt", "copilot"})


def subscription_auth_warnings(
    *, llm_provider: str = "", vision_provider: str = ""
) -> list[str]:
    """Warn when a route signs in with a subscription account (ChatGPT/Copilot).

    Those logins can expire or be invalidated by the provider and force a fresh
    sign-in — most visibly after a restart — whereas API-key providers do not.
    Advisory only; the user's choice is still honored. Returns one message per
    affected role (Vision LLM, Main LLM)."""
    warnings: list[str] = []
    if (vision_provider or "").strip().lower() in _SUBSCRIPTION_AUTH_PROVIDERS:
        warnings.append(
            "Image model uses a subscription login. You may need to sign in again "
            "after restart. For fewer login issues, use an API-key provider."
        )
    if (llm_provider or "").strip().lower() in _SUBSCRIPTION_AUTH_PROVIDERS:
        warnings.append(
            "Chat model uses a subscription login. You may need to sign in again "
            "after restart. For fewer login issues, use an API-key provider."
        )
    return warnings


def _get_tool_schemas(
    prompt: str = "",
    *,
    include_general: bool = True,
    include_screenshot: bool = False,
    allowed_tools: list[str] | None = None,
    pinned_tools: list[str] | None = None,
    unfiltered: bool = False,
) -> list[dict]:
    """Anthropic tool schemas for a query.

    capture_screen is opt-in (include_screenshot) and is never part of the
    general set, so it only appears when the caller explicitly allows the model
    to take screenshots. include_general gates every other built-in tool.

    ``pinned_tools`` are caller-level "always on" tools: they are offered even
    when the per-prompt keyword filter would have dropped them.

    ``unfiltered`` skips the per-prompt keyword filter and returns the full,
    stable built-in set. Multi-turn chat uses this so the tool list (which renders
    at the front of the prompt) stays byte-identical across turns and does not
    invalidate prompt caching — the single-shot query path keeps filtering for
    latency.
    """
    base = (
        _TOOL_REGISTRY.schemas(include_server_tools=True)
        if unfiltered
        else _TOOL_REGISTRY.filtered_schemas(prompt, include_server_tools=True)
    )
    schemas: list[dict] = []
    if include_general:
        schemas = [
            schema
            for schema in base
            if _tool_schema_allowed(schema.get("name", ""), allowed_tools)
        ]
        _append_pinned_tool_schemas(schemas, pinned_tools, allowed_tools)
        _append_local_file_tool_schemas(
            schemas,
            prompt,
            allowed_tools,
            pinned_tools,
            unfiltered=unfiltered,
        )
    if include_screenshot:
        spec = _TOOL_REGISTRY.get_tool("capture_screen")
        if spec is not None:
            schemas.append(spec.anthropic_schema())
    if allowed_tools is not None and _tools_allow(allowed_tools, "memory_search"):
        spec = _TOOL_REGISTRY.get_tool("memory_search")
        if spec is not None:
            schemas.append(spec.anthropic_schema())
    if allowed_tools is not None and _tools_allow(allowed_tools, "memory_save"):
        spec = _TOOL_REGISTRY.get_tool("memory_save")
        if spec is not None:
            schemas.append(spec.anthropic_schema())
    return schemas


def _get_openai_tool_schemas(
    prompt: str = "",
    *,
    include_general: bool = True,
    include_screenshot: bool = False,
    allowed_tools: list[str] | None = None,
    pinned_tools: list[str] | None = None,
) -> list[dict]:
    """OpenAI/Groq function schemas for a query (mirror of _get_tool_schemas).

    capture_screen is opt-in, so it's added back explicitly only when
    include_screenshot is set.
    """
    schemas: list[dict] = []
    if include_general:
        schemas = [
            schema
            for schema in _TOOL_REGISTRY.filtered_openai_schemas(prompt)
            if _tool_schema_allowed(
                ((schema.get("function") or {}).get("name") or ""),
                allowed_tools,
            )
        ]
        _append_pinned_tool_schemas(schemas, pinned_tools, allowed_tools, openai_format=True)
        _append_local_file_tool_schemas(
            schemas,
            prompt,
            allowed_tools,
            pinned_tools,
            openai_format=True,
        )
    if include_screenshot:
        spec = _TOOL_REGISTRY.get_tool("capture_screen")
        if spec is not None:
            schemas.append(spec.openai_schema())
    if allowed_tools is not None and _tools_allow(allowed_tools, "memory_search"):
        spec = _TOOL_REGISTRY.get_tool("memory_search")
        if spec is not None:
            schemas.append(spec.openai_schema())
    if allowed_tools is not None and _tools_allow(allowed_tools, "memory_save"):
        spec = _TOOL_REGISTRY.get_tool("memory_save")
        if spec is not None:
            schemas.append(spec.openai_schema())
    return schemas


def _get_responses_tool_schemas(
    prompt: str = "",
    *,
    include_general: bool = True,
    allowed_tools: list[str] | None = None,
    pinned_tools: list[str] | None = None,
) -> list[dict]:
    """Responses API function schemas for a query."""
    schemas = _get_openai_tool_schemas(
        prompt,
        include_general=include_general,
        include_screenshot=False,
        allowed_tools=allowed_tools,
        pinned_tools=pinned_tools,
    )
    result: list[dict] = []
    for schema in schemas:
        fn = schema.get("function") or {}
        name = str(fn.get("name") or "")
        if not name:
            continue
        result.append({
            "type": "function",
            "name": name,
            "description": str(fn.get("description") or ""),
            "parameters": fn.get("parameters") or {"type": "object", "properties": {}},
        })
    return result


def _tool_schema_allowed(name: str, allowed_tools: list[str] | None) -> bool:
    """Handle tool schema allowed for LLM clients client."""
    if allowed_tools is None:
        return True
    allowed = set(allowed_tools)
    if name == "get_context":
        return bool({"get_context", "get_context.browser", "get_context.documents"} & allowed)
    return name in allowed


def _append_pinned_tool_schemas(
    schemas: list[dict],
    pinned_tools: list[str] | None,
    allowed_tools: list[str] | None,
    *,
    openai_format: bool = False,
) -> None:
    """Add back pinned ("always on") tools the per-prompt keyword filter dropped.

    Mutates *schemas* in place. Opt-in tools (capture_screen, memory_search) are
    skipped — they have dedicated caller settings — and Anthropic-only server
    tools never cross to the OpenAI format.
    """
    if not pinned_tools:
        return
    if openai_format:
        present = {((s.get("function") or {}).get("name") or "") for s in schemas}
    else:
        present = {s.get("name", "") for s in schemas}
    for name in pinned_tools:
        if name in present or not _tool_schema_allowed(name, allowed_tools):
            continue
        spec = _TOOL_REGISTRY.get_tool(name)
        if spec is None or spec.opt_in:
            continue
        if openai_format:
            if spec.server_schema:
                continue
            schemas.append(spec.openai_schema())
        else:
            schemas.append(spec.anthropic_schema())
        present.add(name)


def _append_local_file_tool_schemas(
    schemas: list[dict],
    prompt: str,
    allowed_tools: list[str] | None,
    pinned_tools: list[str] | None,
    *,
    openai_format: bool = False,
    unfiltered: bool = False,
) -> None:
    """Add explicitly allowed local file tools without exposing them by default."""
    if allowed_tools is None:
        return
    allowed = set(allowed_tools)
    pinned = set(pinned_tools or [])
    if openai_format:
        present = {((s.get("function") or {}).get("name") or "") for s in schemas}
    else:
        present = {s.get("name", "") for s in schemas}
    for name in sorted(_LOCAL_FILE_TOOLS):
        if name not in allowed or name in present:
            continue
        if name not in pinned and not unfiltered and not _TOOL_REGISTRY._tool_visible(name, prompt):
            continue
        spec = _TOOL_REGISTRY.get_tool(name)
        if spec is None:
            continue
        schemas.append(spec.openai_schema() if openai_format else spec.anthropic_schema())
        present.add(name)


def _tool_schema_names(schemas: list[dict] | None, *, openai_format: bool = False) -> list[str]:
    """Handle tool schema names for LLM clients client."""
    names: list[str] = []
    for schema in schemas or []:
        if openai_format:
            name = str((schema.get("function") or {}).get("name") or "")
        else:
            name = str(schema.get("name") or "")
        if name:
            names.append(name)
    return names


def _log_offered_model_tools(
    provider: str,
    model: str,
    *,
    allowed_tools: list[str] | None,
    pinned_tools: list[str] | None,
    schemas: list[dict] | None,
    openai_format: bool = False,
) -> None:
    """Log offered model tools."""
    offered = _tool_schema_names(schemas, openai_format=openai_format)
    print(
        f"[llm] tools offered provider={provider or 'unknown'} model={model!r} "
        f"allowed={allowed_tools if allowed_tools is not None else 'all'} "
        f"pinned={pinned_tools or []} offered={offered}",
        flush=True,
    )
    if not offered:
        print(
            "[llm] warning: model tools were requested, but no tool schemas were offered.",
            flush=True,
        )


def _execute_model_tool(name: str, inputs: dict, allowed_tools: list[str] | None = None) -> str:
    """Handle execute model tool for LLM clients client."""
    if allowed_tools is not None:
        allowed = set(allowed_tools)
        if name == "get_context":
            url = str((inputs or {}).get("url") or "").strip()
            needed = "get_context.browser" if url else "get_context.documents"
            if "get_context" not in allowed and needed not in allowed:
                return f"Tool {name!r} is disabled for this context source."
        elif name not in allowed:
            return f"Tool {name!r} is disabled for this caller."
    if name in _LOCAL_FILE_TOOLS:
        return execute_live_file_tool(
            name,
            inputs or {},
            access_mode=_effective_live_file_access_mode(allowed_tools),
            approval_callback=_effective_live_file_approval_callback(),
        )
    return _TOOL_REGISTRY.execute(name, inputs)


def _effective_live_file_access_mode(allowed_tools: list[str] | None = None) -> str:
    """Return the live request's local-file mode, with legacy fallback."""
    explicit = getattr(_LIVE_TOOL_CONTEXT, "file_access_mode", None)
    if explicit:
        return normalize_file_access_mode(explicit)
    allowed = set(allowed_tools or [])
    legacy = str(getattr(config, "TOOL_FILE_MODE", "never") or "never").strip().lower()
    if legacy in {"ask", "auto"} and allowed & {"create_file", "edit_file", "write_file"}:
        return legacy
    if allowed & {"list_files", "read_file"}:
        return "read"
    return normalize_file_access_mode(legacy)


def _tools_allow(allowed_tools: list[str] | None, *names: str) -> bool:
    """Handle tools allow for LLM clients client."""
    if allowed_tools is None:
        return True
    allowed = set(allowed_tools)
    return any(name in allowed for name in names)


def _frontload_context_for_disabled_tools(
    allowed_tools: list[str] | None,
    query: str = "",
) -> str:
    """Handle frontload context for disabled tools for LLM clients client."""
    parts: list[str] = []
    if _tools_allow(allowed_tools, "get_context", "get_context.documents"):
        try:
            doc_text = read_active_document_for_context()
        except Exception as exc:
            print(f"[llm] failed to front-load open document context: {exc}", flush=True)
            doc_text = ""
        if doc_text:
            parts.append("[Open documents]\n" + doc_text)
    if _tools_allow(allowed_tools, "git_status"):
        try:
            status = _execute_git_status({})
        except Exception as exc:
            print(f"[llm] failed to front-load git status context: {exc}", flush=True)
            status = ""
        if status:
            parts.append("[Git status]\n" + status)
    if _tools_allow(allowed_tools, "git_diff"):
        try:
            diff = _execute_git_diff({})
        except Exception as exc:
            print(f"[llm] failed to front-load git diff context: {exc}", flush=True)
            diff = ""
        if diff:
            parts.append("[Git diff]\n" + diff)
    if allowed_tools is not None and _tools_allow(allowed_tools, "memory_search"):
        try:
            from core.memory_store import store

            memory = store.get_manager().retrieve_relevant(query or "") or ""
        except Exception as exc:
            print(f"[llm] failed to front-load memory context: {exc}", flush=True)
            memory = ""
        if memory:
            parts.append(memory)
    return "\n\n".join(parts)


def _append_ambient_context(ambient_context: str, extra: str) -> str:
    """Append ambient context."""
    extra = (extra or "").strip()
    if not extra:
        return ambient_context
    return f"{ambient_context}\n\n---\n{extra}".strip() if ambient_context else extra


def _inject_frontloaded_tool_context(
    ambient_context: str,
    allowed_tools: list[str] | None,
    query: str = "",
) -> str:
    """Handle inject frontloaded tool context for LLM clients client."""
    return _append_ambient_context(
        ambient_context,
        _frontload_context_for_disabled_tools(allowed_tools, query=query),
    )


def _capture_screen_b64(provided_b64: str | None = None) -> str | None:
    """Grab the primary monitor and return it as a base64 PNG, or None on failure."""
    from core.llm_clients.screenshot_tool import resolve_capture_screen_b64

    return resolve_capture_screen_b64(provided_b64)


def _anthropic_vision_model(current_model: str) -> str:
    """Model to use after a screenshot enters the loop.

    Prefer the configured Anthropic vision model; otherwise keep the current
    tool model (Sonnet, which is itself vision-capable).
    """
    if config.VISION_LLM_PROVIDER.strip() == "anthropic" and config.VISION_LLM_MODEL.strip():
        return config.VISION_LLM_MODEL.strip()
    return current_model


def _openai_vision_model(provider: str, current_model: str) -> str:
    """Model to answer with after a screenshot enters an OpenAI/Groq loop.

    Prefer the configured vision model when it lives on the *same* provider
    (we can't swap the bound client mid-loop); otherwise keep the current
    model and hope it is vision-capable (e.g. gpt-4o).
    """
    if config.VISION_LLM_PROVIDER.strip() == provider and config.VISION_LLM_MODEL.strip():
        vision_model = config.VISION_LLM_MODEL.strip()
        if _is_route_cooling(provider, vision_model):
            print(
                f"[llm] screenshot follow-up skipping cooling vision route "
                f"{provider}/{vision_model}; using {current_model}",
                flush=True,
            )
            return current_model
        return vision_model
    return current_model


def _init_keyword_filters() -> None:
    """Handle init keyword filters for LLM clients client."""
    from core.system.paths import TOOL_KEYWORDS_FILE
    _TOOL_REGISTRY.load_keyword_filters(TOOL_KEYWORDS_FILE)
    if not TOOL_KEYWORDS_FILE.exists():
        _TOOL_REGISTRY.save_keyword_filters(TOOL_KEYWORDS_FILE)


_register_builtin_tools()
_init_keyword_filters()


def get_tool_registry() -> ToolRegistry:
    """Return the shared model tool registry. Used by the mod manager."""
    return _TOOL_REGISTRY


def read_document_file(path: str, max_chars: int | None = None) -> str:
    """Read a single local document file and return its plain text (public API)."""
    return _read_document_file(path, max_chars=max_chars)


def read_active_document_for_context_with_debug(active_window: dict | None = None) -> tuple[str, dict]:
    """
    Read all open doc-app windows (foreground and background) and return their
    redacted plain text for proactive injection into the system prompt.
    Multiple documents are separated by per-file headers.
    Returns ("", debug) if no readable documents are found.
    """
    from core.context_fetcher import (
        WindowInfo,
        get_all_open_document_paths,
        get_all_open_document_window_texts,
    )

    active_win = None
    if isinstance(active_window, dict) and (
        active_window.get("title") or active_window.get("name") or active_window.get("window_id")
    ):
        active_win = WindowInfo(
            title=str(active_window.get("title") or active_window.get("name") or ""),
            process_name=str(active_window.get("process_name") or ""),
            pid=int(active_window.get("pid") or 0),
            hwnd=int(active_window.get("window_id") or active_window.get("hwnd") or 0),
        )

    debug: dict = {
        "active_window": active_window or {},
        "paths": [],
        "path_chars": 0,
        "window_labels": [],
        "window_chars": 0,
    }
    paths = get_all_open_document_paths(active_window=active_win)
    debug["paths"] = list(paths)
    if paths:
        text = _read_document_paths(paths)
        debug["path_chars"] = len(text or "")
        print(
            f"[llm] active document paths={len(paths)} paths={paths!r} chars={len(text or '')}",
            flush=True,
        )
        if text:
            return text, debug

    window_texts = get_all_open_document_window_texts(
        max_chars_per_doc=_ambient_document_max_chars(),
        active_window=active_win,
    )
    debug["window_labels"] = [label for label, _text in window_texts]
    debug["window_chars"] = sum(len(text) for _label, text in window_texts)
    print(
        f"[llm] active document window_texts={len(window_texts)} "
        f"labels={debug['window_labels']!r} chars={debug['window_chars']}",
        flush=True,
    )
    if not window_texts:
        return "", debug
    return "\n\n".join(f"[{label}]\n{text}" for label, text in window_texts if text), debug


def read_active_document_for_context() -> str:
    """Read active document for context."""
    text, _debug = read_active_document_for_context_with_debug()
    return text


# ------------------------------------------------------------------
# Singleton clients -” initialised once, reused across all requests
# ------------------------------------------------------------------
_openai_client = None
_anthropic_client = None
_chat_openai_client = None
_chat_anthropic_client = None
_vision_openai_client = None
_vision_anthropic_client = None
_codex_client = None
_chat_codex_client = None

_TEST_IMAGE_BASE64 = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+aF9sAAAAASUVORK5CYII="

# How many times the OpenAI SDK may retry the *same* model inside one call.
_OPENAI_MAX_RETRIES = 1

# Route circuit breaker: when a route returns 429/5xx/no-content, it is parked for this many
# seconds so subsequent turns skip straight to the fallback instead of
# re-probing an exhausted model on every reply. After it expires the route is
# tried again (some retries, not an infinite skip).
_ROUTE_COOLDOWN_SECONDS = 300.0
import time as _time
_route_cooldowns: dict[tuple[str, str], float] = {}
_route_cooldowns_lock = _threading.Lock()
_codex_client_lock = _threading.Lock()


def _route_key(provider: str, model: str) -> tuple[str, str]:
    """Handle route key for LLM clients client."""
    return ((provider or "").lower(), model or "")


def _is_route_cooling(provider: str, model: str) -> bool:
    """Return whether route cooling is true."""
    key = _route_key(provider, model)
    with _route_cooldowns_lock:
        until = _route_cooldowns.get(key)
        if until is None:
            return False
        if _time.time() >= until:
            _route_cooldowns.pop(key, None)
            return False
        return True


def _mark_route_cooling(provider: str, model: str, seconds: float = _ROUTE_COOLDOWN_SECONDS) -> None:
    """Handle mark route cooling for LLM clients client."""
    with _route_cooldowns_lock:
        _route_cooldowns[_route_key(provider, model)] = _time.time() + seconds


def _is_quota_error(exc: Exception) -> bool:
    """True for 429 / rate-limit / quota-exhausted errors worth a cooldown."""
    status = getattr(exc, "status_code", None) or getattr(exc, "status", None)
    if status == 429:
        return True
    name = type(exc).__name__.lower()
    if "ratelimit" in name:
        return True
    text = str(exc).lower()
    return "429" in text or "quota" in text or "rate limit" in text or "rate_limit" in text


def _is_transient_route_error(exc: Exception) -> bool:
    """True for provider-side temporary failures worth trying/skipping fallback."""
    status = getattr(exc, "status_code", None) or getattr(exc, "status", None)
    if status in {429, 500, 502, 503, 504}:
        return True
    text = str(exc).lower()
    return any(
        marker in text
        for marker in (
            "429",
            "500",
            "502",
            "503",
            "504",
            "quota",
            "rate limit",
            "rate_limit",
            "unavailable",
            "high demand",
            "temporarily",
            "try again later",
        )
    )


def _route_failure_summary(
    kind: str,
    attempts: list[tuple[str, str, Exception | str]],
    last_exc: Exception,
) -> RuntimeError:
    """Handle route failure summary for LLM clients client."""
    details = []
    for provider, model, err in attempts:
        details.append(f"{provider}/{model}: {err}")
    joined = "; ".join(details)
    return RuntimeError(f"All {kind} model routes failed. Tried {joined}")


def reset_clients() -> None:
    """Discard all cached API clients so they are rebuilt with the current config."""
    global _openai_client, _anthropic_client, _chat_openai_client, _chat_anthropic_client
    global _vision_openai_client, _vision_anthropic_client, _codex_client, _chat_codex_client
    global _dynamic_anthropic_client_cache
    _openai_client = _anthropic_client = None
    _chat_openai_client = _chat_anthropic_client = None
    _vision_openai_client = _vision_anthropic_client = None
    _codex_client = _chat_codex_client = None
    with _dynamic_client_lock:
        _dynamic_openai_clients.clear()
        _dynamic_anthropic_client_cache = None
    with _route_capabilities_lock:
        _route_capabilities.clear()
    _TOOL_REGISTRY.plugin_dir = Path(config.TOOL_PLUGIN_DIR)
    _TOOL_REGISTRY.refresh()


# ------------------------------------------------------------------
# Codex (ChatGPT subscription) -” custom httpx transport + client
# ------------------------------------------------------------------

class _CodexTransport:
    """
    httpx-compatible BaseTransport that:
      -¢ strips the placeholder API-key authorization header
      -¢ injects ``Authorization: Bearer <access_token>``
      -¢ adds ``ChatGPT-Account-Id`` and ``originator`` headers
    Tokens are fetched (and refreshed) lazily on every request.
    """

    def __init__(self):
        """Initialize the codex transport instance."""
        import httpx
        self._inner = httpx.HTTPTransport()

    def handle_request(self, request):
        """Handle request."""
        import httpx
        from core.auth import chatgpt as chatgpt_auth

        token      = chatgpt_auth.get_valid_access_token()
        account_id = chatgpt_auth.get_account_id()

        # Rebuild raw headers, dropping the dummy API-key auth header
        raw: list[tuple[bytes, bytes]] = [
            (name, value)
            for name, value in request.headers.raw
            if name.lower() != b"authorization"
        ]
        if token:
            raw.append((b"authorization", f"Bearer {token}".encode()))
        if account_id:
            raw.append((b"chatgpt-account-id", account_id.encode()))
        raw.append((b"originator", b"opencode"))

        new_req = httpx.Request(
            method=request.method,
            url=request.url,
            headers=raw,
            content=request.content,
            extensions=request.extensions,
        )
        return self._inner.handle_request(new_req)

    def close(self):
        """Close the wrapped transport."""
        self._inner.close()


def _get_codex_client():
    """Return codex client."""
    global _codex_client
    if _codex_client is None:
        with _codex_client_lock:
            if _codex_client is None:
                with ssl_init_lock():
                    _codex_client = sdk_clients.openai_client(
                        api_key="chatgpt-oauth-dummy",
                        base_url="https://chatgpt.com/backend-api/codex",
                        http_client=sdk_clients.httpx_client(transport=_CodexTransport()),
                    )
    return _codex_client


def _get_chat_codex_client():
    """Returns the same singleton as _get_codex_client() -” same endpoint."""
    return _get_codex_client()


# ------------------------------------------------------------------
# Config sanity checks -” raise early with actionable messages
# ------------------------------------------------------------------

def _log_model_route(kind: str, provider: str, model: str, use_tools: bool = False) -> None:
    """Log model route."""
    import time

    ts = time.strftime("%H:%M:%S")
    tool_note = " tools=on" if use_tools else ""
    print(
        f"[llm {ts}] Route ({kind}): provider={provider or '[unset]'} "
        f"model={model or '[unset]'} credential={_credential_source_for_provider(provider)}{tool_note}"
    )

# Ceiling used when a caller asks for "no app cap" (max_tokens == 0) on a
# provider that, unlike the OpenAI-compatible APIs, requires an explicit limit.
_ANTHROPIC_UNCAPPED_MAX_TOKENS = 8192

# All providers that go through the OpenAI-compatible chat completions API.
_OPENAI_COMPAT_PROVIDER_SET = frozenset({
    "groq", "openai", "google",
    "deepseek", "openrouter", "mistral", "xai", "together", "cerebras", "ollama",
    "custom",
})

_OPENAI_COMPAT_PROVIDERS: dict[str, tuple[str, str]] = {
    # provider → (api_key_attr, base_url)
    "groq":       ("GROQ_API_KEY",       "https://api.groq.com/openai/v1"),
    "google":     ("GOOGLE_API_KEY",     _GOOGLE_OPENAI_BASE_URL),
    "deepseek":   ("DEEPSEEK_API_KEY",   _DEEPSEEK_BASE_URL),
    "openrouter": ("OPENROUTER_API_KEY", _OPENROUTER_BASE_URL),
    "mistral":    ("MISTRAL_API_KEY",    _MISTRAL_BASE_URL),
    "xai":        ("XAI_API_KEY",        _XAI_BASE_URL),
    "together":   ("TOGETHER_API_KEY",   _TOGETHER_BASE_URL),
    "cerebras":   ("CEREBRAS_API_KEY",   _CEREBRAS_BASE_URL),
    "ollama":     (None,                 _OLLAMA_BASE_URL),
}


@dataclass
class RouteCapabilities:
    """Model route capabilities."""
    supports_stream: bool | None = None
    requires_stream: bool | None = None
    supports_tools: bool | None = None
    supports_parallel_tools: bool | None = None
    supports_images: bool | None = None
    supports_json_mode: bool | None = None
    supports_max_output_tokens: bool | None = None
    unsupported_parameters: set[str] = field(default_factory=set)


_route_capabilities: dict[tuple[str, str, str], RouteCapabilities] = {}
_route_capabilities_lock = _threading.Lock()


def _openai_compat_base_url(provider: str) -> str:
    """Handle openai compat base url for LLM clients client."""
    provider = (provider or "").strip().lower()
    if provider == "openai":
        return "https://api.openai.com/v1"
    if provider == "custom":
        return config.CUSTOM_BASE_URL
    item = _OPENAI_COMPAT_PROVIDERS.get(provider)
    return item[1] if item else ""


def _route_endpoint(provider: str) -> str:
    """Handle route endpoint for LLM clients client."""
    provider = (provider or "").strip().lower()
    if provider in _OPENAI_COMPAT_PROVIDER_SET:
        return _openai_compat_base_url(provider).rstrip("/")
    if provider == "chatgpt":
        return "https://chatgpt.com/backend-api/codex"
    if provider == "anthropic":
        return "https://api.anthropic.com"
    if provider == "copilot":
        return "copilot"
    return provider


def _route_capability_key(provider: str, model: str) -> tuple[str, str, str]:
    """Handle route capability key for LLM clients client."""
    return (
        (provider or "").strip().lower(),
        model or "",
        _route_endpoint(provider),
    )


def _get_route_capabilities(provider: str, model: str) -> RouteCapabilities:
    """Return route capabilities."""
    key = _route_capability_key(provider, model)
    with _route_capabilities_lock:
        cap = _route_capabilities.get(key)
        if cap is None:
            cap = RouteCapabilities()
            _route_capabilities[key] = cap
        return cap


def _update_route_capabilities(provider: str, model: str, **facts) -> None:
    """Update route capabilities."""
    cap = _get_route_capabilities(provider, model)
    for name, value in facts.items():
        if hasattr(cap, name):
            setattr(cap, name, value)


def _mark_unsupported_parameter(provider: str, model: str, name: str) -> None:
    """Handle mark unsupported parameter for LLM clients client."""
    if not name:
        return
    cap = _get_route_capabilities(provider, model)
    cap.unsupported_parameters.add(name)
    if name in {"max_output_tokens", "max_tokens"}:
        cap.supports_max_output_tokens = False
    elif name == "tools":
        cap.supports_tools = False
    elif name == "parallel_tool_calls":
        cap.supports_parallel_tools = False
    elif name == "response_format":
        cap.supports_json_mode = False


def _recover_openai_compat_kwargs(
    provider: str, model: str, kwargs: dict, exc: Exception
) -> dict | None:
    """Return kwargs adjusted to satisfy a 400 about a specific field, or None.

    The general "comply with whatever the endpoint rejects" path: rename
    ``max_tokens`` to ``max_completion_tokens`` when the model demands it, or drop
    any unsupported parameter/value we can name (temperature, top_p, etc.). The
    fix is recorded on the route so later calls skip the retry. Returns None when
    the error is not a recoverable field problem (caller re-raises)."""
    _record_route_error_capabilities(provider, model, exc)
    lowered = str(exc).lower()
    if "max_tokens" in kwargs and "max_completion_tokens" in lowered:
        new = dict(kwargs)
        new["max_completion_tokens"] = new.pop("max_tokens")
        _mark_unsupported_parameter(provider, model, "max_tokens")
        print(f"[llm] {provider}/{model}: max_tokens -> max_completion_tokens", flush=True)
        return new
    name = _unsupported_parameter_name(exc)
    if name and name in kwargs:
        _mark_unsupported_parameter(provider, model, name)
        new = dict(kwargs)
        new.pop(name, None)
        print(f"[llm] {provider}/{model}: dropping unsupported field {name!r}", flush=True)
        return new
    return None


def _stream_openai_compat_plain(
    provider: str, model: str, kwargs: dict
) -> Generator[str, None, None]:
    """Stream a no-tool OpenAI-compatible completion, self-healing the request.

    On a 400 about a specific field it drops or renames that field and retries,
    so one code path works across models with different format rules. Honors the
    learned per-route capabilities, the macOS non-streaming safe mode, and the
    streaming -> non-streaming fallback. Used by the chat and rewrite routes."""
    kwargs = dict(kwargs)
    for unsupported in list(_get_route_capabilities(provider, model).unsupported_parameters):
        kwargs.pop(unsupported, None)

    if not kwargs.get("stream", True):
        while True:
            try:
                text = _openai_compat_stdlib_completion_text(provider, kwargs)
                break
            except Exception as exc:
                adjusted = _recover_openai_compat_kwargs(provider, model, kwargs, exc)
                if adjusted is None:
                    raise
                kwargs = adjusted
        if text:
            yield text
        return

    client = _dynamic_openai_client(provider)
    while True:
        produced = False
        try:
            with client.chat.completions.create(**kwargs) as stream:
                for chunk in stream:
                    if not chunk.choices:
                        continue
                    delta = chunk.choices[0].delta.content
                    if delta:
                        produced = True
                        yield delta
            return
        except Exception as exc:
            if produced:
                raise  # already streamed output; don't retry and double-emit
            adjusted = _recover_openai_compat_kwargs(provider, model, kwargs, exc)
            if adjusted is not None:
                kwargs = adjusted
                continue
            if kwargs.get("stream") and _stream_mode_error(exc):
                if not _stream_false_required_error(exc):
                    _update_route_capabilities(provider, model, supports_stream=False, requires_stream=False)
                kwargs = dict(kwargs)
                kwargs["stream"] = False
                print("[llm] OpenAI-compatible stream rejected; retrying non-streaming", flush=True)
                while True:
                    try:
                        response = client.chat.completions.create(**kwargs)
                        break
                    except Exception as create_exc:
                        adjusted = _recover_openai_compat_kwargs(provider, model, kwargs, create_exc)
                        if adjusted is None:
                            raise
                        kwargs = adjusted
                text = _openai_compat_message_text(response)
                if text:
                    yield text
                return
            raise


def _openai_compat_api_key(provider: str) -> str:
    """Handle openai compat api key for LLM clients client."""
    provider = (provider or "").strip().lower()
    if provider == "openai":
        return config.OPENAI_API_KEY
    if provider == "custom":
        return config.CUSTOM_API_KEY or "no-key"
    item = _OPENAI_COMPAT_PROVIDERS.get(provider)
    if not item:
        return ""
    key_attr, _base_url = item
    return getattr(config, key_attr) if key_attr else "ollama"


def _use_macos_openai_compat_non_streaming(provider: str) -> bool:
    """Avoid macOS native crashes in OpenAI-compatible streaming paths."""
    return not macos_safety.openai_compat_streaming_enabled(provider)


_openai_compat_stdlib_ssl_context = None
_openai_compat_stdlib_ssl_context_lock = _threading.Lock()


def _get_openai_compat_stdlib_ssl_context():
    """Return a cached certifi-backed SSL context for urllib compatibility calls."""
    global _openai_compat_stdlib_ssl_context
    with _openai_compat_stdlib_ssl_context_lock:
        if _openai_compat_stdlib_ssl_context is None:
            import certifi

            _openai_compat_stdlib_ssl_context = _ssl.create_default_context(cafile=certifi.where())
        return _openai_compat_stdlib_ssl_context


def _openai_compat_message_text(response) -> str:
    """Handle openai compat message text for LLM clients client."""
    choice = response.choices[0] if getattr(response, "choices", None) else None
    if choice is None:
        return ""
    message = getattr(choice, "message", None)
    content = getattr(message, "content", "") if message is not None else ""
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts: list[str] = []
        for item in content:
            if isinstance(item, str):
                parts.append(item)
            elif isinstance(item, dict):
                parts.append(str(item.get("text") or ""))
            else:
                parts.append(str(getattr(item, "text", "") or ""))
        return "".join(parts).strip()
    return str(content or "").strip()


def _openai_compat_message_tool_calls(response) -> dict[int, dict]:
    """Handle openai compat message tool calls for LLM clients client."""
    choice = response.choices[0] if getattr(response, "choices", None) else None
    message = getattr(choice, "message", None) if choice is not None else None
    tool_calls = getattr(message, "tool_calls", None) or []
    acc: dict[int, dict] = {}
    for i, tc in enumerate(tool_calls):
        fn = getattr(tc, "function", None)
        acc[i] = {
            "id": getattr(tc, "id", "") or f"tool_call_{i}",
            "name": getattr(fn, "name", "") if fn is not None else "",
            "arguments": getattr(fn, "arguments", "") if fn is not None else "",
        }
    return acc


def _response_output_text(response) -> str:
    """Handle response output text for LLM clients client."""
    if isinstance(response, dict):
        text = response.get("output_text", "")
        output = response.get("output", []) or []
    else:
        text = getattr(response, "output_text", "")
        output = getattr(response, "output", []) or []
    if isinstance(text, str) and text:
        return text
    parts: list[str] = []
    for item in output:
        content_items = item.get("content", []) if isinstance(item, dict) else getattr(item, "content", [])
        for content in content_items or []:
            value = content.get("text", "") if isinstance(content, dict) else getattr(content, "text", "")
            if value:
                parts.append(str(value))
    return "".join(parts).strip()


def _response_id(response) -> str:
    """Return a Responses API id from object or dict response data."""
    if isinstance(response, dict):
        return str(response.get("id") or "")
    return str(getattr(response, "id", "") or "")


def _response_output_items(response) -> list:
    """Return Responses API output items from object or dict response data."""
    if isinstance(response, dict):
        output = response.get("output", []) or []
    else:
        output = getattr(response, "output", []) or []
    return list(output)


def _response_function_calls(response) -> list[dict[str, str]]:
    """Extract function_call items from a Responses API response."""
    calls: list[dict[str, str]] = []
    for item in _response_output_items(response):
        if isinstance(item, dict):
            item_type = str(item.get("type") or "")
            item_id = str(item.get("id") or "")
            name = str(item.get("name") or "")
            arguments = str(item.get("arguments") or "")
            call_id = str(item.get("call_id") or item_id or "")
        else:
            item_type = str(getattr(item, "type", "") or "")
            item_id = str(getattr(item, "id", "") or "")
            name = str(getattr(item, "name", "") or "")
            arguments = str(getattr(item, "arguments", "") or "")
            call_id = str(getattr(item, "call_id", "") or item_id or "")
        if item_type == "function_call" and name:
            calls.append({"id": item_id, "name": name, "arguments": arguments, "call_id": call_id})
    return calls


def _event_value(obj, name: str, default=None):
    """Read an event/object field from either a dict or SDK object."""
    if isinstance(obj, dict):
        return obj.get(name, default)
    return getattr(obj, name, default)


def _normalized_response_item(item) -> dict:
    """Convert a streamed Responses output item to the shape create() returns."""
    if isinstance(item, dict):
        return dict(item)
    normalized: dict = {}
    for field_name in ("id", "type", "call_id", "name", "arguments", "content", "status", "role"):
        value = getattr(item, field_name, None)
        if value is not None:
            normalized[field_name] = value
    return normalized


def _responses_stream_to_response(client, kwargs: dict, *, provider: str, model: str):
    """Stream a Responses request and reconstruct a minimal response object."""
    current = dict(kwargs)
    while True:
        response_id = ""
        output_text: list[str] = []
        output_items: dict[str, dict] = {}
        completed_response = None
        try:
            with client.responses.stream(**_responses_stream_kwargs(current)) as stream:
                _update_route_capabilities(provider, model, supports_stream=True, requires_stream=True)
                for event in stream:
                    event_type = str(_event_value(event, "type", "") or "")
                    response = _event_value(event, "response", None)
                    if response is not None:
                        response_id = response_id or _response_id(response)
                    if event_type == "response.completed":
                        completed_response = response
                        continue
                    if event_type == "response.output_text.delta":
                        delta = str(_event_value(event, "delta", "") or "")
                        if delta:
                            output_text.append(delta)
                        continue
                    if event_type in {"response.output_item.added", "response.output_item.done"}:
                        item = _event_value(event, "item", None)
                        if item is not None:
                            normalized = _normalized_response_item(item)
                            output_index = _event_value(event, "output_index", None)
                            key = str(
                                output_index
                                if output_index not in (None, "")
                                else normalized.get("id")
                                or normalized.get("call_id")
                                or len(output_items)
                            )
                            existing = output_items.get(key, {})
                            merged = {**existing, **{k: v for k, v in normalized.items() if v not in (None, "")}}
                            output_items[key] = merged
                        continue
                    if event_type in {
                        "response.function_call_arguments.delta",
                        "response.function_call_arguments.done",
                    }:
                        output_index = _event_value(event, "output_index", None)
                        key = str(
                            output_index
                            if output_index not in (None, "")
                            else _event_value(event, "item_id", "")
                            or _event_value(event, "output_index", "")
                            or len(output_items)
                        )
                        item = output_items.setdefault(key, {"type": "function_call"})
                        for field_name in ("call_id", "name"):
                            value = _event_value(event, field_name, None)
                            if value:
                                item[field_name] = value
                        arguments = str(_event_value(event, "arguments", "") or "")
                        if event_type.endswith(".delta"):
                            item["arguments"] = str(item.get("arguments") or "") + str(
                                _event_value(event, "delta", "") or ""
                            )
                        elif arguments:
                            item["arguments"] = arguments
        except Exception as exc:
            _record_route_error_capabilities(provider, model, exc)
            if _store_must_be_false_error(exc) and current.get("store") is True:
                current = dict(current)
                current["store"] = False
                print("[llm] Responses stream requires store=false; retrying without stored response", flush=True)
                continue
            retry_kwargs = _without_unsupported_parameter(current, exc)
            if retry_kwargs is None:
                raise
            name = _unsupported_parameter_name(exc)
            _mark_unsupported_parameter(provider, model, name)
            if name == "tools":
                raise
            print("[llm] Responses stream rejected unsupported parameter; retrying without it", flush=True)
            current = retry_kwargs
            continue

        if completed_response is not None:
            completed_output = _response_output_items(completed_response)
            if completed_output or _response_output_text(completed_response):
                return completed_response
        return {
            "id": response_id,
            "output_text": "".join(output_text),
            "output": list(output_items.values()),
        }


def _responses_create_with_retries(client, kwargs: dict, *, provider: str, model: str):
    """Create a Responses API response, retrying unsupported parameters."""
    current = dict(kwargs)
    while True:
        try:
            return client.responses.create(**current)
        except Exception as exc:
            _record_route_error_capabilities(provider, model, exc)
            if _store_must_be_false_error(exc) and current.get("store") is True:
                current = dict(current)
                current["store"] = False
                print("[llm] Responses create requires store=false; retrying without stored response", flush=True)
                continue
            if _requires_stream_error(exc):
                _update_route_capabilities(provider, model, supports_stream=True, requires_stream=True)
                return _responses_stream_to_response(client, current, provider=provider, model=model)
            retry_kwargs = _without_unsupported_parameter(current, exc)
            if retry_kwargs is None:
                raise
            name = _unsupported_parameter_name(exc)
            _mark_unsupported_parameter(provider, model, name)
            if name == "tools":
                raise
            print("[llm] Responses create rejected unsupported parameter; retrying without it", flush=True)
            current = retry_kwargs


def _responses_tool_loop_kwargs(kwargs: dict) -> dict:
    """Keep Responses tool-call turns available for later function outputs."""
    current = dict(kwargs)
    # Responses function_call_output items are matched against prior tool calls
    # via previous_response_id. Some routes cannot do that if the prior response
    # was created with store=false, so tool loops must be stateful.
    current["store"] = True
    return current


def _responses_stream_kwargs(kwargs: dict) -> dict:
    """Return kwargs for Responses streaming with an explicit body stream flag."""
    current = dict(kwargs)
    extra_body = dict(current.get("extra_body") or {})
    extra_body["stream"] = True
    current["extra_body"] = extra_body
    return current


def _no_matching_tool_call_error(exc: Exception) -> bool:
    """Return True when a Responses route rejects a tool output linkage."""
    text = str(exc).lower()
    return "no tool call found" in text and "function call output" in text


def _store_must_be_false_error(exc: Exception) -> bool:
    """Return True when a Responses route rejects stored responses."""
    text = str(exc).lower()
    return "store" in text and "false" in text and any(
        marker in text
        for marker in (
            "must be",
            "must set",
            "set to",
            "should be",
            "only supports",
        )
    )


def _stateless_tool_output_input(calls: list[dict[str, str]], tool_outputs: list[dict]) -> list[dict]:
    """Build a fallback input that carries function calls and outputs together."""
    items: list[dict] = []
    for call, output in zip(calls, tool_outputs):
        call_item = {
            "type": "function_call",
            "call_id": call["call_id"],
            "name": call["name"],
            "arguments": call.get("arguments") or "{}",
        }
        if call.get("id"):
            call_item["id"] = call["id"]
        items.append(call_item)
        items.append(output)
    return items


def _run_responses_tool_loop(
    client,
    kwargs: dict,
    *,
    provider: str,
    model: str,
    allowed_tools: list[str] | None = None,
    max_rounds: int = 3,
) -> Generator[str, None, None]:
    """Run a non-streaming Responses API function-call loop."""
    followup_instructions = kwargs.get("instructions")
    response = _responses_create_with_retries(
        client,
        _responses_tool_loop_kwargs(kwargs),
        provider=provider,
        model=model,
    )
    _update_route_capabilities(provider, model, supports_tools=True)
    for _round in range(max_rounds):
        calls = _response_function_calls(response)
        if not calls:
            text = _response_output_text(response)
            if text:
                yield text
            return
        tool_outputs = []
        for call in calls:
            try:
                inputs = _stdlib_json.loads(call["arguments"] or "{}")
                if not isinstance(inputs, dict):
                    inputs = {}
            except Exception:
                inputs = {}
            result = _execute_model_tool(call["name"], inputs, allowed_tools=allowed_tools)
            tool_outputs.append({
                "type": "function_call_output",
                "call_id": call["call_id"],
                "output": result,
            })
        followup_kwargs = {
            "model": model,
            "input": tool_outputs,
            "previous_response_id": _response_id(response),
            "store": True,
        }
        if followup_instructions:
            followup_kwargs["instructions"] = followup_instructions
        try:
            response = _responses_create_with_retries(
                client,
                followup_kwargs,
                provider=provider,
                model=model,
            )
        except Exception as exc:
            if not _no_matching_tool_call_error(exc):
                raise
            fallback_kwargs = {
                "model": model,
                "input": _stateless_tool_output_input(calls, tool_outputs),
                "store": False,
            }
            if followup_instructions:
                fallback_kwargs["instructions"] = followup_instructions
            response = _responses_create_with_retries(
                client,
                fallback_kwargs,
                provider=provider,
                model=model,
            )
    text = _response_output_text(response)
    if text:
        yield text


def _stream_mode_error(exc: Exception) -> bool:
    """Stream mode error."""
    status = getattr(exc, "status_code", None) or getattr(exc, "status", None)
    text = str(exc).lower()
    if "stream" not in text and "streaming" not in text:
        return False
    return status in {400, 404, 409, 422, None} or any(
        marker in text
        for marker in (
            "unsupported",
            "not supported",
            "must be",
            "is not true",
            "isn't true",
            "isnt true",
            "required",
            "invalid",
            "unrecognized",
        )
    )


def _requires_stream_error(exc: Exception) -> bool:
    """Handle requires stream error for LLM clients client."""
    text = str(exc).lower()
    return any(
        marker in text
        for marker in (
            "stream must be true",
            "stream must be set to true",
            "stream is not true",
            "stream isn't true",
            "stream isnt true",
            "stream=true",
            "requires stream",
            "requires streaming",
            "must be streamed",
        )
    )


def _streaming_not_supported_error(exc: Exception) -> bool:
    """Handle streaming not supported error for LLM clients client."""
    text = str(exc).lower()
    if _requires_stream_error(exc):
        return False
    return any(
        marker in text
        for marker in (
            "streaming is not supported",
            "stream is not supported",
            "streaming not supported",
            "stream not supported",
            "streaming unsupported",
            "unsupported stream",
        )
    )


def _stream_false_required_error(exc: Exception) -> bool:
    """Return whether this request, not the whole route, must be non-streaming."""
    text = str(exc).lower()
    return any(
        marker in text
        for marker in (
            "stream must be false",
            "stream=false",
            "not compatible with streaming",
            "not compatible with stream",
        )
    )


def _tools_not_supported_error(exc: Exception) -> bool:
    """Handle tools not supported error for LLM clients client."""
    text = str(exc).lower()
    if "tool" not in text and "function" not in text:
        return False
    return any(
        marker in text
        for marker in (
            "unsupported",
            "not supported",
            "unknown parameter",
            "unrecognized",
            "invalid parameter",
            "not allowed",
        )
    )


def _json_mode_not_supported_error(exc: Exception) -> bool:
    """Handle json mode not supported error for LLM clients client."""
    text = str(exc).lower()
    return (
        "response_format" in text
        or "json mode" in text
        or ("json" in text and any(marker in text for marker in ("unsupported", "not supported")))
    )


def _image_not_supported_error(exc: Exception) -> bool:
    """Handle image not supported error for LLM clients client."""
    text = str(exc).lower()
    return (
        any(marker in text for marker in ("image", "vision", "input_image", "image_url"))
        and any(marker in text for marker in ("unsupported", "not supported", "invalid", "not accept"))
    )


def _record_route_error_capabilities(provider: str, model: str, exc: Exception) -> None:
    """Record route error capabilities."""
    name = _unsupported_parameter_name(exc)
    if name:
        _mark_unsupported_parameter(provider, model, name)
    if _requires_stream_error(exc):
        _update_route_capabilities(provider, model, supports_stream=True, requires_stream=True)
    elif _streaming_not_supported_error(exc):
        _update_route_capabilities(provider, model, supports_stream=False, requires_stream=False)
    if (not name or name == "tools") and _tools_not_supported_error(exc):
        _update_route_capabilities(provider, model, supports_tools=False)
    if _json_mode_not_supported_error(exc):
        _update_route_capabilities(provider, model, supports_json_mode=False)
    if _image_not_supported_error(exc):
        _update_route_capabilities(provider, model, supports_images=False)


def _unsupported_parameter_name(exc: Exception) -> str:
    """Handle unsupported parameter name for LLM clients client."""
    text = str(exc)
    lowered = text.lower()
    # "Unsupported value: 'temperature' does not support 0.5 with this model. Only
    # the default (1) value is supported." — GPT-5-family / reasoning models reject
    # any non-default sampling value. Treat the named field as unsupported so the
    # route drops it and falls back to the model's default, instead of failing.
    if "unsupported value" in lowered and (
        "only the default" in lowered or "does not support" in lowered
    ):
        rest = text[lowered.find("unsupported value"):]
        import re as _re
        m = _re.search(r"['\"]([a-zA-Z0-9_]+)['\"]", rest)
        if m:
            return m.group(1)
    if "unsupported parameter" not in lowered and "unknown parameter" not in lowered:
        return ""
    for marker in ("parameter:", "parameter", "param:"):
        idx = lowered.find(marker)
        if idx >= 0:
            raw = text[idx + len(marker):].strip()
            if raw.startswith(("'", '"')):
                quote = raw[0]
                end = raw.find(quote, 1)
                if end > 1:
                    return raw[1:end]
            return raw.split()[0].strip(" .,:;'\"{}[]")
    return ""


def _without_unsupported_parameter(kwargs: dict, exc: Exception) -> dict | None:
    """Handle without unsupported parameter for LLM clients client."""
    name = _unsupported_parameter_name(exc)
    if not name or name not in kwargs:
        return None
    retry_kwargs = dict(kwargs)
    retry_kwargs.pop(name, None)
    return retry_kwargs


def _response_stream_text(
    client,
    kwargs: dict,
    *,
    provider: str = "",
    model: str = "",
) -> Generator[str, None, None]:
    """Handle response stream text for LLM clients client."""
    kwargs = dict(kwargs)
    cap = _get_route_capabilities(provider, model) if provider or model else None
    if cap is not None:
        for name in list(cap.unsupported_parameters):
            kwargs.pop(name, None)
        if cap.requires_stream is False or cap.supports_stream is False:
            try:
                response = client.responses.create(**kwargs)
            except Exception as create_exc:
                _record_route_error_capabilities(provider, model, create_exc)
                if _requires_stream_error(create_exc):
                    _update_route_capabilities(provider, model, supports_stream=True, requires_stream=True)
                    yield from _response_stream_text(client, kwargs, provider=provider, model=model)
                    return
                retry_kwargs = _without_unsupported_parameter(kwargs, create_exc)
                if retry_kwargs is None:
                    raise
                _mark_unsupported_parameter(provider, model, _unsupported_parameter_name(create_exc))
                print("[llm] Responses create rejected unsupported parameter; retrying without it", flush=True)
                response = client.responses.create(**retry_kwargs)
            text = _response_output_text(response)
            if text:
                yield text
            return
    try:
        with client.responses.stream(**_responses_stream_kwargs(kwargs)) as stream:
            if cap is not None:
                cap.supports_stream = True
            for event in stream:
                if getattr(event, "type", "") == "response.output_text.delta":
                    delta = getattr(event, "delta", "")
                    if delta:
                        yield delta
    except Exception as exc:
        if provider or model:
            _record_route_error_capabilities(provider, model, exc)
        retry_kwargs = _without_unsupported_parameter(kwargs, exc)
        if retry_kwargs is not None:
            if provider or model:
                _mark_unsupported_parameter(provider, model, _unsupported_parameter_name(exc))
            print("[llm] Responses stream rejected unsupported parameter; retrying without it", flush=True)
            yield from _response_stream_text(client, retry_kwargs, provider=provider, model=model)
            return
        if not _stream_mode_error(exc):
            raise
        if provider or model:
            if not _stream_false_required_error(exc):
                _update_route_capabilities(provider, model, supports_stream=False, requires_stream=False)
        print("[llm] Responses stream rejected; retrying non-streaming", flush=True)
        try:
            response = client.responses.create(**kwargs)
        except Exception as create_exc:
            if provider or model:
                _record_route_error_capabilities(provider, model, create_exc)
            if _requires_stream_error(create_exc):
                if provider or model:
                    _update_route_capabilities(provider, model, supports_stream=True, requires_stream=True)
                yield from _response_stream_text(client, kwargs, provider=provider, model=model)
                return
            retry_kwargs = _without_unsupported_parameter(kwargs, create_exc)
            if retry_kwargs is None:
                raise
            if provider or model:
                _mark_unsupported_parameter(provider, model, _unsupported_parameter_name(create_exc))
            print("[llm] Responses create rejected unsupported parameter; retrying without it", flush=True)
            response = client.responses.create(**retry_kwargs)
        text = _response_output_text(response)
        if text:
            yield text


def _openai_compat_stdlib_completion_text(provider: str, kwargs: dict) -> str:
    """Blocking OpenAI-compatible completion without importing provider SDKs."""
    def _request_completion() -> str:
        """Handle request completion for LLM clients client."""
        base_url = _openai_compat_base_url(provider).rstrip("/")
        if not base_url:
            raise ValueError(f"No OpenAI-compatible base URL configured for {provider!r}")
        api_key = _openai_compat_api_key(provider)
        url = f"{base_url}/chat/completions"
        body = _stdlib_json.dumps(kwargs).encode("utf-8")
        headers = {"Content-Type": "application/json"}
        if api_key and provider != "ollama":
            headers["Authorization"] = f"Bearer {api_key}"
        request = _urllib_request.Request(url, data=body, headers=headers, method="POST")
        opener = _urllib_request.build_opener(
            _urllib_request.ProxyHandler({}),
            _urllib_request.HTTPSHandler(context=_get_openai_compat_stdlib_ssl_context()),
        )
        try:
            with opener.open(request, timeout=60) as response:
                raw = response.read()
                encoding = response.headers.get("Content-Encoding", "")
        except _urllib_error.HTTPError as exc:
            raw = exc.read()
            encoding = exc.headers.get("Content-Encoding", "") if exc.headers else ""
            if "gzip" in encoding.lower() and raw:
                raw = gzip.decompress(raw)
            message = raw.decode("utf-8", errors="replace") if raw else str(exc)
            raise RuntimeError(f"Error code: {exc.code} - {message}") from exc
        if "gzip" in encoding.lower() and raw:
            raw = gzip.decompress(raw)
        payload = _stdlib_json.loads(raw.decode("utf-8"))
        choices = payload.get("choices") or []
        if not choices:
            return ""
        message = choices[0].get("message") or {}
        content = message.get("content") or ""
        if isinstance(content, list):
            parts: list[str] = []
            for item in content:
                if isinstance(item, str):
                    parts.append(item)
                elif isinstance(item, dict):
                    parts.append(str(item.get("text") or ""))
                else:
                    parts.append(str(getattr(item, "text", "") or ""))
            return "".join(parts).strip()
        return str(content).strip()

    with native_init_lock():
        return _request_completion()


# Per-provider cache for the route/fallback clients. Building these on every
# query rebuilt an SSL context each time; worse, the build raced with the TTS
# client's SSL build (separate thread, same query) and segfaulted on macOS.
# Cache + serialize construction so the SSL context is built once, never
# concurrently. Cleared by reset_clients() when settings change.
_dynamic_openai_clients: dict[str, object] = {}
_dynamic_anthropic_client_cache = None
_dynamic_client_lock = _threading.Lock()


def _build_dynamic_openai_client(provider: str):
    # max_retries=1: allow one quick same-model retry for a transient blip, but
    # not the SDK default (2) which slowly re-hammers an exhausted model every
    # turn with backoff before the fallback chain (which switches models) is
    # ever reached. Quota cooldown (see _stream_with_fallbacks) handles the rest.
    """Build dynamic openai client."""
    if provider in _OPENAI_COMPAT_PROVIDERS:
        key_attr, base_url = _OPENAI_COMPAT_PROVIDERS[provider]
        api_key = getattr(config, key_attr) if key_attr else "ollama"
        return sdk_clients.openai_client(api_key=api_key or "no-key", base_url=base_url, max_retries=_OPENAI_MAX_RETRIES)
    if provider == "custom":
        return sdk_clients.openai_client(api_key=config.CUSTOM_API_KEY or "no-key", base_url=config.CUSTOM_BASE_URL, max_retries=_OPENAI_MAX_RETRIES)
    return sdk_clients.openai_client(api_key=config.OPENAI_API_KEY, max_retries=_OPENAI_MAX_RETRIES)


def _dynamic_openai_client(provider: str):
    """Handle dynamic openai client for LLM clients client."""
    with _dynamic_client_lock:
        client = _dynamic_openai_clients.get(provider)
        if client is None:
            with ssl_init_lock():
                client = _build_dynamic_openai_client(provider)
            _dynamic_openai_clients[provider] = client
        return client


def _dynamic_anthropic_client():
    """Handle dynamic anthropic client for LLM clients client."""
    global _dynamic_anthropic_client_cache
    with _dynamic_client_lock:
        if _dynamic_anthropic_client_cache is None:
            with ssl_init_lock():
                _dynamic_anthropic_client_cache = sdk_clients.anthropic_client(api_key=config.ANTHROPIC_API_KEY)
        return _dynamic_anthropic_client_cache


def prewarm() -> None:
    """Build the primary query-route client once at startup.

    Runs sequentially with tts.prewarm() on the startup thread, so the LLM and
    TTS SSL contexts are built one after another here instead of racing on the
    first query (the macOS segfault). Also removes the handshake from the first
    query. Best-effort: a missing key or unreachable provider must not crash
    startup — the real query path handles those.
    """
    try:
        provider = (config.LLM_PROVIDER or "").strip().lower()
        if provider == "anthropic":
            _dynamic_anthropic_client()
        elif provider in _OPENAI_COMPAT_PROVIDER_SET:
            _dynamic_openai_client(provider)
        # chatgpt/copilot use the Codex transport, which builds no SSL context here.
    except Exception:
        pass


def list_models(provider: str, *, api_key: str = "", base_url: str = "") -> list[str]:
    """Return available model ids for *provider*, fetched live from its API.

    Raises on any failure — missing credential, network error, or a provider
    that has no public listing endpoint (chatgpt/copilot). Callers should treat
    a raise as "fall back to the curated list". Optional ``api_key``/``base_url``
    let the settings dialog fetch with not-yet-saved field values.
    """
    provider = (provider or "").lower()
    if provider in ("chatgpt", "copilot"):
        raise NotImplementedError(f"{provider} does not support model listing")

    if provider == "anthropic":
        key = api_key or config.ANTHROPIC_API_KEY
        if not key:
            raise ValueError("No Anthropic API key configured")
        with ssl_init_lock():
            client = sdk_clients.anthropic_client(api_key=key)
        resp = client.models.list(limit=1000)
        ids = [m.id for m in resp.data]
    else:
        if provider == "custom":
            url = base_url or config.CUSTOM_BASE_URL
            if not url:
                raise ValueError("No custom base URL configured")
            with ssl_init_lock():
                client = sdk_clients.openai_client(api_key=api_key or config.CUSTOM_API_KEY or "no-key", base_url=url)
        elif provider in _OPENAI_COMPAT_PROVIDERS:
            key_attr, default_base = _OPENAI_COMPAT_PROVIDERS[provider]
            key = api_key or (getattr(config, key_attr) if key_attr else "ollama")
            if key_attr and not key:
                raise ValueError(f"No API key configured for {provider}")
            with ssl_init_lock():
                client = sdk_clients.openai_client(api_key=key or "no-key", base_url=base_url or default_base)
        elif provider == "openai":
            key = api_key or config.OPENAI_API_KEY
            if not key:
                raise ValueError("No OpenAI API key configured")
            with ssl_init_lock():
                client = sdk_clients.openai_client(api_key=key)
        else:
            raise ValueError(f"Unknown provider: {provider}")
        resp = client.models.list()
        ids = [_normalize_model_for_provider(provider, m.id) for m in resp.data]

    if not ids:
        raise ValueError("Provider returned no models")
    return sorted(ids)


_CHATGPT_SUPPORTED_MODELS = {
    "gpt-5.5", "gpt-5.4", "gpt-5.4-mini",
    "gpt-5.3-codex", "gpt-5.3-codex-spark", "gpt-5.2",
}


def _check_llm_config() -> None:
    """Check llm config."""
    if not config.LLM_MODEL:
        raise ValueError(
            "LLM_MODEL is not set in .env. "
            "Add LLM_MODEL=<model> (e.g. llama3-8b-8192 for Groq)."
        )
    if config.LLM_PROVIDER.lower() == "chatgpt":
        from core.auth import chatgpt as chatgpt_auth
        if not chatgpt_auth.get_tokens():
            raise ValueError(
                "LLM_PROVIDER is set to 'chatgpt' but you are not logged in. "
                "Open Settings -> LLM and sign in with your ChatGPT account."
            )
        if config.LLM_MODEL.lower() not in _CHATGPT_SUPPORTED_MODELS:
            raise ValueError(
                f"Model '{config.LLM_MODEL}' is not supported by the ChatGPT Codex endpoint. "
                f"Use one of: {', '.join(sorted(_CHATGPT_SUPPORTED_MODELS))}"
            )
        return
    if config.LLM_PROVIDER.lower() == "copilot":
        from core.auth import copilot_auth
        if not copilot_auth.get_token():
            raise ValueError(
                "LLM_PROVIDER is set to 'copilot' but no GitHub Copilot token is stored. "
                "Open Settings -> LLM and save a GitHub Copilot token."
            )
        return
    if config.LLM_PROVIDER.lower() == "custom" and not config.CUSTOM_BASE_URL:
        raise ValueError(
            "LLM_PROVIDER is set to 'custom' but CUSTOM_BASE_URL is not configured. "
            "Add the base URL in Settings → LLM → Custom provider."
        )
    if config.LLM_PROVIDER.lower() == "ollama":
        return
    if not _api_key_for(config.LLM_PROVIDER):
        raise ValueError(
            f"API key for LLM_PROVIDER='{config.LLM_PROVIDER}' is not configured. "
            "Add the matching key in Settings so it can be stored in the OS keychain."
        )


def _check_chat_llm_config() -> None:
    """Check chat llm config."""
    if not config.CHAT_LLM_MODEL:
        raise ValueError(
            "CHAT_LLM_MODEL is not set in .env. "
            "Add CHAT_LLM_MODEL=<model> or leave unset to inherit LLM_MODEL."
        )
    if config.CHAT_LLM_PROVIDER.lower() == "chatgpt":
        from core.auth import chatgpt as chatgpt_auth
        if not chatgpt_auth.get_tokens():
            raise ValueError(
                "CHAT_LLM_PROVIDER is set to 'chatgpt' but you are not logged in. "
                "Open Settings -> LLM and sign in with your ChatGPT account."
            )
        if config.CHAT_LLM_MODEL.lower() not in _CHATGPT_SUPPORTED_MODELS:
            raise ValueError(
                f"Model '{config.CHAT_LLM_MODEL}' is not supported by the ChatGPT Codex endpoint. "
                f"Use one of: {', '.join(sorted(_CHATGPT_SUPPORTED_MODELS))}"
            )
        return
    if config.CHAT_LLM_PROVIDER.lower() == "copilot":
        from core.auth import copilot_auth
        if not copilot_auth.get_token():
            raise ValueError(
                "CHAT_LLM_PROVIDER is set to 'copilot' but no GitHub Copilot token is stored. "
                "Open Settings -> LLM and save a GitHub Copilot token."
            )
        return
    if config.CHAT_LLM_PROVIDER.lower() == "custom" and not config.CUSTOM_BASE_URL:
        raise ValueError(
            "CHAT_LLM_PROVIDER is set to 'custom' but CUSTOM_BASE_URL is not configured. "
            "Add the base URL in Settings → LLM → Custom provider."
        )
    if config.CHAT_LLM_PROVIDER.lower() == "ollama":
        return
    if not _api_key_for(config.CHAT_LLM_PROVIDER):
        raise ValueError(
            f"API key for CHAT_LLM_PROVIDER='{config.CHAT_LLM_PROVIDER}' is not configured. "
            "Add the matching key in Settings so it can be stored in the OS keychain."
        )


def _check_vision_config() -> None:
    """Check vision config."""
    if not config.VISION_LLM_MODEL:
        raise ValueError(
            "VISION_LLM_MODEL is not set in .env. "
            "Screen-snip queries require a vision-capable model. "
            "Example: VISION_LLM_PROVIDER=anthropic  VISION_LLM_MODEL=claude-opus-4-5"
        )
    if not config.VISION_LLM_PROVIDER:
        raise ValueError(
            "VISION_LLM_PROVIDER is not set in .env. "
            "Add VISION_LLM_PROVIDER=anthropic (or openai or chatgpt) alongside VISION_LLM_MODEL."
        )
    if config.VISION_LLM_PROVIDER.lower() == "chatgpt":
        from core.auth import chatgpt as chatgpt_auth
        if not chatgpt_auth.get_tokens():
            raise ValueError(
                "VISION_LLM_PROVIDER is set to 'chatgpt' but you are not logged in. "
                "Open Settings -> LLM and sign in with your ChatGPT account."
            )
        return
    if config.VISION_LLM_PROVIDER.lower() == "custom" and not config.CUSTOM_BASE_URL:
        raise ValueError(
            "VISION_LLM_PROVIDER is set to 'custom' but CUSTOM_BASE_URL is not configured. "
            "Add the base URL in Settings → LLM → Custom provider."
        )
    if config.VISION_LLM_PROVIDER.lower() == "ollama":
        return
    if not _api_key_for(config.VISION_LLM_PROVIDER):
        raise ValueError(
            f"API key for VISION_LLM_PROVIDER='{config.VISION_LLM_PROVIDER}' is not configured. "
            "Add the matching key in Settings so it can be stored in the OS keychain."
        )


def _check_route_config(provider: str, model: str, route_name: str) -> None:
    """Check route config."""
    if not provider or not model:
        raise ValueError(f"{route_name} route is missing provider or model.")
    if provider == "chatgpt":
        from core.auth import chatgpt as chatgpt_auth
        if not chatgpt_auth.get_tokens():
            raise ValueError(f"{route_name} route uses chatgpt but you are not logged in.")
        if model.lower() not in _CHATGPT_SUPPORTED_MODELS:
            raise ValueError(
                f"Model '{model}' is not supported by the ChatGPT Codex endpoint. "
                f"Use one of: {', '.join(sorted(_CHATGPT_SUPPORTED_MODELS))}"
            )
        return
    if provider == "copilot":
        from core.auth import copilot_auth
        if not copilot_auth.get_token():
            raise ValueError(f"{route_name} route uses copilot but no GitHub Copilot token is stored.")
        return
    if provider == "custom":
        if not config.CUSTOM_BASE_URL:
            raise ValueError(f"{route_name} route uses 'custom' but CUSTOM_BASE_URL is not set.")
        return
    if provider == "ollama":
        return   # local, no key required
    if not _api_key_for(provider):
        raise ValueError(f"{route_name} route uses {provider!r}, but its API key is not configured.")


def _check_route_config_with_credentials(
    provider: str,
    model: str,
    route_name: str,
    *,
    anthropic_api_key: str = "",
    custom_base_url: str = "",
    compat_keys: dict[str, str] | None = None,
) -> None:
    """Check route config with credentials."""
    if not provider or not model:
        raise ValueError(f"{route_name} route is missing provider or model.")
    if provider == "chatgpt":
        from core.auth import chatgpt as chatgpt_auth
        if not chatgpt_auth.get_tokens():
            raise ValueError(f"{route_name} route uses chatgpt but you are not logged in.")
        if model.lower() not in _CHATGPT_SUPPORTED_MODELS:
            raise ValueError(
                f"Model '{model}' is not supported by the ChatGPT Codex endpoint. "
                f"Use one of: {', '.join(sorted(_CHATGPT_SUPPORTED_MODELS))}"
            )
        return
    if provider == "copilot":
        from core.auth import copilot_auth
        if not copilot_auth.get_token():
            raise ValueError(f"{route_name} route uses copilot but no GitHub Copilot token is stored.")
        return
    if provider == "custom":
        if not (custom_base_url or config.CUSTOM_BASE_URL):
            raise ValueError(f"{route_name} route uses 'custom' but CUSTOM_BASE_URL is not set.")
        return
    if provider == "ollama":
        return   # no key required
    if provider == "anthropic":
        if not anthropic_api_key:
            raise ValueError(f"{route_name} route uses anthropic, but its API key is not configured.")
        return
    if provider in _OPENAI_COMPAT_PROVIDER_SET:
        api_key = (
            compat_keys.get(provider, "")
            if compat_keys is not None
            else _api_key_for(provider)
        )
        if not api_key:
            raise ValueError(f"{route_name} route uses {provider!r}, but its API key is not configured.")
        return
    if not _api_key_for(provider):
        raise ValueError(f"{route_name} route uses {provider!r}, but its API key is not configured.")


def _get_openai_client():
    """Return openai client."""
    global _openai_client
    if _openai_client is None:
        with ssl_init_lock():
            if config.LLM_PROVIDER.lower() == "groq":
                _openai_client = sdk_clients.openai_client(
                    api_key=config.GROQ_API_KEY,
                    base_url="https://api.groq.com/openai/v1",
                )
            else:
                _openai_client = sdk_clients.openai_client(api_key=config.OPENAI_API_KEY)
    return _openai_client


def _get_chat_openai_client():
    """Returns the same singleton as _get_openai_client() when providers match."""
    if config.CHAT_LLM_PROVIDER.lower() == config.LLM_PROVIDER.lower():
        return _get_openai_client()
    global _chat_openai_client
    if _chat_openai_client is None:
        with ssl_init_lock():
            if config.CHAT_LLM_PROVIDER.lower() == "groq":
                _chat_openai_client = sdk_clients.openai_client(
                    api_key=config.GROQ_API_KEY,
                    base_url="https://api.groq.com/openai/v1",
                )
            else:
                _chat_openai_client = sdk_clients.openai_client(api_key=config.OPENAI_API_KEY)
    return _chat_openai_client


def _get_anthropic_client():
    """Return anthropic client."""
    global _anthropic_client
    if _anthropic_client is None:
        with ssl_init_lock():
            _anthropic_client = sdk_clients.anthropic_client(api_key=config.ANTHROPIC_API_KEY)
    return _anthropic_client


def _get_chat_anthropic_client():
    """Return chat anthropic client."""
    if config.CHAT_LLM_PROVIDER.lower() == config.LLM_PROVIDER.lower():
        return _get_anthropic_client()
    global _chat_anthropic_client
    if _chat_anthropic_client is None:
        with ssl_init_lock():
            _chat_anthropic_client = sdk_clients.anthropic_client(api_key=config.ANTHROPIC_API_KEY)
    return _chat_anthropic_client


def _get_vision_openai_client():
    """Return vision openai client."""
    global _vision_openai_client
    if _vision_openai_client is None:
        with ssl_init_lock():
            if config.VISION_LLM_PROVIDER.lower() == "groq":
                _vision_openai_client = sdk_clients.openai_client(
                    api_key=config.GROQ_API_KEY,
                    base_url="https://api.groq.com/openai/v1",
                )
            elif config.VISION_LLM_PROVIDER.lower() == "google":
                _vision_openai_client = sdk_clients.openai_client(
                    api_key=config.GOOGLE_API_KEY,
                    base_url=_GOOGLE_OPENAI_BASE_URL,
                )
            else:
                _vision_openai_client = sdk_clients.openai_client(api_key=config.OPENAI_API_KEY)
    return _vision_openai_client


def _get_vision_anthropic_client():
    """Return vision anthropic client."""
    global _vision_anthropic_client
    if _vision_anthropic_client is None:
        with ssl_init_lock():
            _vision_anthropic_client = sdk_clients.anthropic_client(api_key=config.ANTHROPIC_API_KEY)
    return _vision_anthropic_client


def _run_openai_compat_probe(client, *, provider: str, model: str, messages: list) -> None:
    """Run openai compat probe."""
    cap = _get_route_capabilities(provider, model)
    kwargs = {
        "model": model,
        "messages": messages,
        "stream": cap.supports_stream is not False and not _use_macos_openai_compat_non_streaming(provider),
        "max_tokens": 8,
    }
    for unsupported in list(cap.unsupported_parameters):
        kwargs.pop(unsupported, None)
    if not kwargs["stream"]:
        client.chat.completions.create(**kwargs)
        return
    try:
        with client.chat.completions.create(**kwargs) as stream:
            _update_route_capabilities(provider, model, supports_stream=True)
            for _chunk in stream:
                break
    except Exception as exc:
        _record_route_error_capabilities(provider, model, exc)
        retry_kwargs = _without_unsupported_parameter(kwargs, exc)
        if retry_kwargs is not None:
            _mark_unsupported_parameter(provider, model, _unsupported_parameter_name(exc))
            print("[llm] route probe rejected unsupported parameter; retrying without it", flush=True)
            with client.chat.completions.create(**retry_kwargs) as stream:
                for _chunk in stream:
                    break
            return
        if not _stream_mode_error(exc):
            raise
        retry_kwargs = dict(kwargs)
        retry_kwargs["stream"] = False
        _update_route_capabilities(provider, model, supports_stream=False, requires_stream=False)
        print("[llm] route probe stream rejected; retrying non-streaming", flush=True)
        client.chat.completions.create(**retry_kwargs)


def _probe_openai_compat_route(provider: str, model: str, image_base64: str | None = None) -> None:
    """Handle probe openai compat route for LLM clients client."""
    client = _dynamic_openai_client(provider)
    _run_openai_compat_probe(
        client,
        provider=provider,
        model=_normalize_model_for_provider(provider, model),
        messages=_build_openai_messages("Reply with OK.", image_base64),
    )


def _probe_openai_compat_route_with_credentials(
    provider: str,
    model: str,
    *,
    api_key: str,
    base_url: str = "",
    image_base64: str | None = None,
) -> None:
    """Handle probe openai compat route with credentials for LLM clients client."""
    with ssl_init_lock():
        if provider == "custom":
            client = sdk_clients.openai_client(api_key=api_key or "no-key", base_url=base_url or config.CUSTOM_BASE_URL)
        elif provider in _OPENAI_COMPAT_PROVIDERS:
            _key_attr, provider_base_url = _OPENAI_COMPAT_PROVIDERS[provider]
            client = sdk_clients.openai_client(api_key=api_key or "no-key", base_url=provider_base_url)
        else:
            client = sdk_clients.openai_client(api_key=api_key)
    _run_openai_compat_probe(
        client,
        provider=provider,
        model=_normalize_model_for_provider(provider, model),
        messages=_build_openai_messages("Reply with OK.", image_base64),
    )


def _probe_anthropic_route(model: str, image_base64: str | None = None) -> None:
    """Handle probe anthropic route for LLM clients client."""
    client = _dynamic_anthropic_client()
    if image_base64:
        content = [
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/png",
                    "data": image_base64,
                },
            },
            {"type": "text", "text": "Reply with OK."},
        ]
    else:
        content = "Reply with OK."
    client.messages.create(
        model=model,
        max_tokens=8,
        messages=[{"role": "user", "content": content}],
    )


def _probe_anthropic_route_with_api_key(model: str, api_key: str, image_base64: str | None = None) -> None:
    """Handle probe anthropic route with api key for LLM clients client."""
    with ssl_init_lock():
        client = sdk_clients.anthropic_client(api_key=api_key)
    if image_base64:
        content = [
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/png",
                    "data": image_base64,
                },
            },
            {"type": "text", "text": "Reply with OK."},
        ]
    else:
        content = "Reply with OK."
    client.messages.create(
        model=model,
        max_tokens=8,
        messages=[{"role": "user", "content": content}],
    )


def _probe_chatgpt_route(model: str, image_base64: str | None = None) -> None:
    """Handle probe chatgpt route for LLM clients client."""
    client = _get_codex_client()
    if image_base64:
        content = [
            {"type": "input_text", "text": "Reply with OK."},
            {"type": "input_image", "image_url": f"data:image/png;base64,{image_base64}"},
        ]
    else:
        content = [{"type": "input_text", "text": "Reply with OK."}]
    kwargs = {
        "model": model,
        "input": [{"type": "message", "role": "user", "content": content}],
        "instructions": "Return exactly OK.",
        "store": False,
        "max_output_tokens": 8,
    }
    try:
        client.responses.create(**kwargs)
        _update_route_capabilities("chatgpt", model, requires_stream=False, supports_max_output_tokens=True)
    except Exception as exc:
        _record_route_error_capabilities("chatgpt", model, exc)
        retry_kwargs = _without_unsupported_parameter(kwargs, exc)
        if retry_kwargs is not None:
            _mark_unsupported_parameter("chatgpt", model, _unsupported_parameter_name(exc))
            print("[llm] ChatGPT route probe create rejected unsupported parameter; retrying without it", flush=True)
            client.responses.create(**retry_kwargs)
            return
        if not _stream_mode_error(exc):
            raise
        _update_route_capabilities("chatgpt", model, supports_stream=True, requires_stream=True)
        print("[llm] ChatGPT route probe create rejected; retrying streaming", flush=True)
        try:
            with client.responses.stream(**_responses_stream_kwargs(kwargs)) as stream:
                for _event in stream:
                    break
        except Exception as stream_exc:
            _record_route_error_capabilities("chatgpt", model, stream_exc)
            retry_kwargs = _without_unsupported_parameter(kwargs, stream_exc)
            if retry_kwargs is None:
                raise
            _mark_unsupported_parameter("chatgpt", model, _unsupported_parameter_name(stream_exc))
            print("[llm] ChatGPT route probe stream rejected unsupported parameter; retrying without it", flush=True)
            with client.responses.stream(**_responses_stream_kwargs(retry_kwargs)) as stream:
                for _event in stream:
                    break


def _probe_copilot_route(model: str) -> None:
    """Handle probe copilot route for LLM clients client."""
    from core.auth import copilot_client

    text = copilot_client.ask(
        "Reply with OK.",
        model,
        system="Return exactly OK.",
        allow_tools=False,
    )
    if not text.strip():
        raise RuntimeError("Copilot returned an empty response.")


def test_route_connection(
    provider: str,
    model: str,
    route_name: str = "LLM",
    *,
    image: bool = False,
    anthropic_api_key: str | None = None,
    custom_base_url: str | None = None,
    compat_keys: dict[str, str] | None = None,
) -> tuple[bool, str]:
    """Test a provider/model route.

    compat_keys: map of provider name → API key for OpenAI-compat providers,
                 used when the dialog passes freshly-typed (not-yet-saved) keys.
    """
    provider = (provider or "").strip().lower()
    model = (model or "").strip()
    try:
        explicit_credentials = bool(compat_keys or anthropic_api_key is not None or custom_base_url is not None)
        if explicit_credentials:
            _check_route_config_with_credentials(
                provider,
                model,
                route_name,
                anthropic_api_key=anthropic_api_key or "",
                custom_base_url=custom_base_url or "",
                compat_keys=compat_keys or {},
            )
        else:
            _check_route_config(provider, model, route_name)

        def _probe_compat(image_b64=None):
            """Handle probe compat for local."""
            if explicit_credentials and compat_keys and provider in compat_keys:
                _probe_openai_compat_route_with_credentials(
                    provider, model,
                    api_key=compat_keys.get(provider, ""),
                    base_url=custom_base_url or "",
                    image_base64=image_b64,
                )
            else:
                _probe_openai_compat_route(provider, model, image_base64=image_b64)

        if image:
            image_base64 = _TEST_IMAGE_BASE64
            if provider in _OPENAI_COMPAT_PROVIDER_SET:
                _probe_compat(image_base64)
            elif provider == "anthropic":
                if explicit_credentials:
                    _probe_anthropic_route_with_api_key(model, anthropic_api_key or "", image_base64=image_base64)
                else:
                    _probe_anthropic_route(model, image_base64=image_base64)
            elif provider == "chatgpt":
                _probe_chatgpt_route(model, image_base64=image_base64)
            else:
                raise ValueError(f"Unknown vision provider: {provider}")
            return True, f"{route_name} vision route OK: {provider} / {model}"

        if provider in _OPENAI_COMPAT_PROVIDER_SET:
            _probe_compat()
        elif provider == "anthropic":
            if explicit_credentials:
                _probe_anthropic_route_with_api_key(model, anthropic_api_key or "")
            else:
                _probe_anthropic_route(model)
        elif provider == "chatgpt":
            _probe_chatgpt_route(model)
        elif provider == "copilot":
            _probe_copilot_route(model)
        else:
            raise ValueError(f"Unknown LLM provider: {provider}")
        return True, f"{route_name} route OK: {provider} / {model}"
    except Exception as exc:
        return False, f"{route_name} test failed: {exc}"


def stream_response(
    user_message: str,
    image_base64: str | None = None,
    ambient_context: str = "",
    memory_context: str = "",
    use_tools: bool = False,
    allowed_tools: list[str] | None = None,
    pinned_tools: list[str] | None = None,
    allow_screenshot_tool: bool = False,
    screenshot_tool_b64: str | None = None,
    route_provider: str | None = None,
    route_model: str | None = None,
    route_fallbacks: str | None = None,
    max_tokens: int | None = None,
    temperature: float | None = None,
    json_mode: bool = False,
    history: list[dict] | None = None,
) -> Generator[str, None, None]:
    """
    Stream a response from the configured LLM.

    ``history`` is an optional list of prior {role, content} turns (the active
    conversation) replayed before the current message so hotkey/voice prompts
    can continue a thread with full model context.

    When image_base64 is provided, uses VISION_LLM_PROVIDER/MODEL.
    Otherwise uses LLM_PROVIDER/MODEL.

    Args:
        user_message:     The user's query text.
        image_base64:     Optional base64-encoded PNG for vision input.
        ambient_context:  Plain-text context block (active window, clipboard,
                          focused element) -” injected into system prompt.
        memory_context:   Pre-formatted LTM facts + STM session summary from
                          core.memory -” injected into system prompt before the
                          ambient context block.
        use_tools:        If True and provider is Anthropic, expose
                  web_search + get_context tools so Claude can
                  pull extra context when it decides to. The model
                  must use the actual tool call interface rather than
                  describing or simulating tool calls in text.
                  Ignored for Groq/OpenAI providers and vision calls.
        allow_screenshot_tool: If True (Anthropic, non-vision), additionally
                  expose the capture_screen tool so the model can grab a
                  screenshot on demand. Independent of use_tools.
        route_provider:   Optional provider override for non-vision calls.
        route_model:      Optional model override for non-vision calls.
        route_fallbacks:  Optional provider:model fallback string for overrides.
        max_tokens:       Optional response budget override for callers that
                          need structured or longer output.
        temperature:      Optional sampling override for callers that need
                          deterministic structured output.

    Yields:
        Text chunks as they arrive from the API.
    """
    if image_base64:
        log_event("llm.user_message", f"User message (vision): {user_message!r}", message_preview=user_message)
    else:
        log_event("llm.user_message", f"User message: {user_message!r}", message_preview=user_message)

    if ambient_context:
        _log_context(
            "ambient snapshot -” captured at hotkey/voice trigger, injected into system prompt",
            ambient_context,
        )

    if image_base64:
        candidates = _route_candidates(
            config.VISION_LLM_PROVIDER,
            config.VISION_LLM_MODEL,
            config.VISION_LLM_FALLBACKS,
        )
        yield from _stream_with_fallbacks(
            "vision",
            candidates,
            lambda provider, model: _stream_single_response_route(
                provider,
                model,
                user_message,
                image_base64,
                ambient_context,
                memory_context,
                use_tools=False,
                route_name="VISION_LLM",
                max_tokens=max_tokens,
                temperature=temperature,
            ),
        )
    else:
        provider = (route_provider or config.LLM_PROVIDER).strip()
        model = (route_model or config.LLM_MODEL).strip()
        fallback_raw = config.LLM_FALLBACKS if route_fallbacks is None else route_fallbacks
        candidates = _route_candidates(provider, model, fallback_raw)
        yield from _stream_with_fallbacks(
            "query",
            candidates,
            lambda provider, model: _stream_single_response_route(
                provider,
                model,
                user_message,
                None,
                ambient_context,
                memory_context,
                use_tools=use_tools,
                allowed_tools=allowed_tools,
                pinned_tools=pinned_tools,
                allow_screenshot_tool=allow_screenshot_tool,
                screenshot_tool_b64=screenshot_tool_b64,
                route_name="LLM",
                max_tokens=max_tokens,
                temperature=temperature,
                json_mode=json_mode,
                history=history,
            ),
        )


def _stream_with_fallbacks(
    kind: str,
    candidates: list[tuple[str, str]],
    factory,
) -> Generator[str, None, None]:
    """Stream with fallbacks."""
    import time
    # Try routes not in quota cooldown first; keep cooling ones as last resort so
    # a fully-throttled chain still attempts something rather than failing cold.
    ready = [c for c in candidates if not _is_route_cooling(*c)]
    cooling = [c for c in candidates if _is_route_cooling(*c)]
    if cooling and not ready:
        tried = "; ".join(f"{provider}/{model}" for provider, model in cooling)
        raise RuntimeError(
            f"All {kind} model routes are temporarily cooling down after recent "
            f"provider failures. Routes: {tried}"
        )
    ordered = ready + cooling
    if not ordered:
        ordered = candidates
    last_exc: Exception | None = None
    attempts: list[tuple[str, str, Exception | str]] = []
    for idx, (provider, model) in enumerate(ordered):
        emitted = False
        route_started = time.monotonic()
        try:
            for chunk in factory(provider, model):
                emitted = True
                yield chunk
            if emitted:
                elapsed = time.monotonic() - route_started
                log_event(
                    "llm.route_complete",
                    f"Route ({kind}) {provider}/{model} completed in {elapsed:.1f}s",
                    kind=kind,
                    provider=provider,
                    model=model,
                    elapsed_seconds=elapsed,
                )
                return
            # Model returned HTTP 200 but zero content chunks — treat as failure.
            last_exc = ValueError(f"Route ({kind}) {provider}/{model} returned no content")
            attempts.append((provider, model, "returned no content"))
            if not _is_route_cooling(provider, model):
                _mark_route_cooling(provider, model)
            elapsed = time.monotonic() - route_started
            if idx < len(ordered) - 1:
                log_event(
                    "llm.route_empty_fallback",
                    f"Route ({kind}) {provider}/{model} returned no content after {elapsed:.1f}s; cooling down {_ROUTE_COOLDOWN_SECONDS:.0f}s and trying fallback",
                    kind=kind,
                    provider=provider,
                    model=model,
                    elapsed_seconds=elapsed,
                )
                continue
            log_event(
                "llm.route_empty_final",
                f"Route ({kind}) {provider}/{model} returned no content after {elapsed:.1f}s; cooling down {_ROUTE_COOLDOWN_SECONDS:.0f}s; no fallback left",
                kind=kind,
                provider=provider,
                model=model,
                elapsed_seconds=elapsed,
            )
        except Exception as exc:
            _record_route_error_capabilities(provider, model, exc)
            last_exc = exc
            attempts.append((provider, model, exc))
            elapsed = time.monotonic() - route_started
            if emitted:
                log_event(
                    "llm.route_streaming_failure",
                    f"Route ({kind}) {provider}/{model} failed after streaming for {elapsed:.1f}s; not falling back: {exc}",
                    kind=kind,
                    provider=provider,
                    model=model,
                    elapsed_seconds=elapsed,
                    error=str(exc),
                )
                raise
            if not _is_route_cooling(provider, model) and _is_transient_route_error(exc):
                _mark_route_cooling(provider, model)
                log_event(
                    "llm.route_transient_fallback",
                    f"Route ({kind}) {provider}/{model} hit transient provider error after {elapsed:.1f}s; cooling down {_ROUTE_COOLDOWN_SECONDS:.0f}s and using fallback",
                    kind=kind,
                    provider=provider,
                    model=model,
                    elapsed_seconds=elapsed,
                    error=str(exc),
                )
            if idx < len(ordered) - 1:
                log_event(
                    "llm.route_failure_fallback",
                    f"Route ({kind}) {provider}/{model} failed before output after {elapsed:.1f}s; trying fallback: {exc}",
                    kind=kind,
                    provider=provider,
                    model=model,
                    elapsed_seconds=elapsed,
                    error=str(exc),
                )
                continue
            log_event(
                "llm.route_failure_final",
                f"Route ({kind}) {provider}/{model} failed after {elapsed:.1f}s; no fallback left: {exc}",
                kind=kind,
                provider=provider,
                model=model,
                elapsed_seconds=elapsed,
                error=str(exc),
            )
    if last_exc:
        raise _route_failure_summary(kind, attempts, last_exc) from last_exc
    raise ValueError(f"No {kind} model routes configured.")


def _stream_single_response_route(
    provider: str,
    model: str,
    user_message: str,
    image_base64: str | None,
    ambient_context: str,
    memory_context: str,
    use_tools: bool,
    route_name: str,
    allowed_tools: list[str] | None = None,
    pinned_tools: list[str] | None = None,
    allow_screenshot_tool: bool = False,
    screenshot_tool_b64: str | None = None,
    max_tokens: int | None = None,
    temperature: float | None = None,
    json_mode: bool = False,
    history: list[dict] | None = None,
    system_prompt: str | None = None,
    route_kind: str = "query",
) -> Generator[str, None, None]:
    """Stream single response route."""
    _check_route_config(provider, model, route_name)
    model = _normalize_model_for_provider(provider, model)
    # capture_screen also needs the Anthropic tool loop, so it counts as "tools".
    anthropic_tools = use_tools or allow_screenshot_tool
    # Use the user's model for tool calls. TOOL_LLM_MODEL is an optional override
    # (empty by default) for users who want a different model when tools are active
    # — we no longer force a hardcoded model.
    effective_model = (
        config.TOOL_LLM_MODEL
        if (provider == "anthropic" and anthropic_tools and config.TOOL_LLM_MODEL.strip())
        else model
    )
    tools_for_log = anthropic_tools and not image_base64
    if provider in _OPENAI_COMPAT_PROVIDER_SET and not macos_safety.openai_compat_tools_enabled():
        tools_for_log = False
    _log_model_route("vision" if image_base64 else route_kind, provider, effective_model, use_tools=tools_for_log)
    if image_base64:
        if provider in _OPENAI_COMPAT_PROVIDER_SET:
            client = None if _use_macos_openai_compat_non_streaming(provider) else _dynamic_openai_client(provider)
            yield from _stream_openai_compat(
                user_message,
                image_base64,
                model,
                client,
                ambient_context,
                memory_context,
                provider=provider,
                max_tokens=max_tokens,
                temperature=temperature,
            )
        elif provider == "anthropic":
            yield from _stream_anthropic(
                user_message,
                image_base64,
                model,
                _dynamic_anthropic_client(),
                ambient_context,
                memory_context,
                max_tokens=max_tokens,
                temperature=temperature,
            )
        elif provider == "chatgpt":
            yield from _stream_codex_vision(
                user_message,
                image_base64,
                model,
                _get_codex_client(),
                ambient_context,
                memory_context,
            )
        else:
            raise ValueError(f"Unknown vision provider: {provider}")
        return
    if provider in _OPENAI_COMPAT_PROVIDER_SET:
        client = None if _use_macos_openai_compat_non_streaming(provider) else _dynamic_openai_client(provider)
        yield from _stream_openai_compat(
            user_message,
            None,
            model,
            client,
            ambient_context,
            memory_context,
            use_tools=use_tools,
            allowed_tools=allowed_tools,
            pinned_tools=pinned_tools,
            allow_screenshot_tool=allow_screenshot_tool,
            screenshot_tool_b64=screenshot_tool_b64,
            provider=provider,
            max_tokens=max_tokens,
            temperature=temperature,
            json_mode=json_mode,
            history=history,
            system_prompt=system_prompt,
        )
    elif provider == "anthropic":
        yield from _stream_anthropic(
            user_message,
            None,
            effective_model,
            _dynamic_anthropic_client(),
            ambient_context,
            memory_context,
            use_tools,
            allowed_tools=allowed_tools,
            pinned_tools=pinned_tools,
            allow_screenshot_tool=allow_screenshot_tool,
            screenshot_tool_b64=screenshot_tool_b64,
            max_tokens=max_tokens,
            temperature=temperature,
            history=history,
            system_prompt=system_prompt,
        )
    elif provider == "chatgpt":
        yield from _stream_codex(
            user_message,
            model,
            _get_codex_client(),
            ambient_context,
            memory_context,
            use_tools,
            allowed_tools=allowed_tools,
            pinned_tools=pinned_tools,
            history=history,
            system_prompt=system_prompt,
        )
    elif provider == "copilot":
        yield from _stream_copilot(
            user_message,
            model,
            ambient_context,
            memory_context,
            use_tools,
            allowed_tools=allowed_tools,
            history=history,
            system_prompt=system_prompt,
        )
    else:
        raise ValueError(f"Unknown LLM provider: {provider}")


def _stream_copilot(
    user_message: str,
    model: str,
    ambient_context: str = "",
    memory_context: str = "",
    use_tools: bool = False,
    allowed_tools: list[str] | None = None,
    history: list[dict] | None = None,
    system_prompt: str | None = None,
) -> Generator[str, None, None]:
    """Stream copilot."""
    from core.auth import copilot_client

    if use_tools:
        _update_route_capabilities("copilot", model, supports_tools=False)
        ambient_context = _inject_frontloaded_tool_context(
            ambient_context,
            allowed_tools,
            query=user_message,
        )
    parts = []
    if system_prompt:
        parts.append(system_prompt)
    if memory_context:
        parts.append(memory_context)
    if ambient_context:
        parts.append("Context:\n" + ambient_context)
    system = "\n\n".join(parts)
    history_prefix = _codex_history_prefix(_sanitize_history(history))
    prompt = (history_prefix + user_message) if history_prefix else user_message
    yield from copilot_client.stream(
        prompt,
        model,
        system=system,
        session_id="wisp-chat" if history else None,
        allow_tools=False,
    )


# ------------------------------------------------------------------
# OpenAI / Groq (OpenAI-compatible)
# ------------------------------------------------------------------

def _stream_openai_compat(
    user_message: str,
    image_base64: str | None,
    model: str,
    client,
    ambient_context: str = "",
    memory_context: str = "",
    use_tools: bool = False,
    allowed_tools: list[str] | None = None,
    pinned_tools: list[str] | None = None,
    allow_screenshot_tool: bool = False,
    screenshot_tool_b64: str | None = None,
    provider: str = "",
    max_tokens: int | None = None,
    temperature: float | None = None,
    json_mode: bool = False,
    history: list[dict] | None = None,
    system_prompt: str | None = None,
) -> Generator[str, None, None]:
    """Stream openai compat."""
    import json as _json

    provider = (provider or "").strip().lower()
    cap = _get_route_capabilities(provider, model)
    if image_base64 and cap.supports_images is False:
        raise RuntimeError(f"{provider}/{model} does not support image input.")
    tools_requested = (use_tools or allow_screenshot_tool) and not image_base64 and not json_mode
    tools_allowed = macos_safety.openai_compat_tools_enabled() and cap.supports_tools is not False
    nonstream_safe_mode = _use_macos_openai_compat_non_streaming(provider)
    if tools_requested and not tools_allowed:
        reason = "macOS safe mode" if not macos_safety.openai_compat_tools_enabled() else "route capability cache"
        print(f"[llm] OpenAI-compatible live tools disabled by {reason}", flush=True)
        ambient_context = _inject_frontloaded_tool_context(
            ambient_context,
            allowed_tools,
            query=user_message,
        )
    if tools_requested and tools_allowed and nonstream_safe_mode:
        print(
            "[llm] OpenAI-compatible non-streaming safe mode cannot run live "
            "tool loops; front-loading supported local context instead.",
            flush=True,
        )
        _update_route_capabilities(provider, model, supports_tools=False)
        ambient_context = _inject_frontloaded_tool_context(
            ambient_context,
            allowed_tools,
            query=user_message,
        )
        tools_requested = False
        tools_allowed = False
        use_tools = False
        allow_screenshot_tool = False
    messages = _build_openai_messages(
        user_message,
        image_base64,
        ambient_context,
        memory_context,
        history,
        system_prompt=system_prompt,
    )
    expose_screenshot = tools_allowed and allow_screenshot_tool and not image_base64 and not json_mode
    if expose_screenshot and messages and messages[0].get("role") == "system":
        messages[0]["content"] = _with_screenshot_note(messages[0]["content"], True)
    # JSON mode (used by the agent protocol) forces a single JSON object as the
    # whole response. Native context tools would let the model emit a function
    # call instead of that JSON, so they are mutually exclusive: the agent embeds
    # its own tool calls inside the protocol JSON and does not need them here.
    tools = (
        _get_openai_tool_schemas(
            user_message,
            include_general=use_tools,
            include_screenshot=allow_screenshot_tool,
            allowed_tools=allowed_tools,
            pinned_tools=pinned_tools,
        )
        if tools_requested and tools_allowed
        else None
    )
    if tools_requested and tools_allowed:
        _log_offered_model_tools(
            provider or "openai-compatible",
            model,
            allowed_tools=allowed_tools,
            pinned_tools=pinned_tools,
            schemas=tools,
            openai_format=True,
        )
    if tools and messages and messages[0].get("role") == "system":
        messages[0]["content"] = _with_tools_note(messages[0]["content"], True)
        messages[0]["content"] = _with_memory_search_note(messages[0]["content"], allowed_tools)
        messages[0]["content"] = _with_memory_save_note(messages[0]["content"], allowed_tools)

    kwargs: dict = {
        "model": model,
        "messages": messages,
        "stream": True,
    }
    _apply_sampling(kwargs, model, 0.5 if temperature is None else temperature)
    # max_tokens == 0 means "no app-imposed cap": omit the field so the provider
    # uses its own per-model maximum. None means "unspecified" -> a safe default.
    if max_tokens != 0:
        _apply_max_output(kwargs, model, max_tokens or 2048)
    if tools:
        kwargs["tools"] = tools
        if cap.supports_parallel_tools is True:
            kwargs["parallel_tool_calls"] = True
    if json_mode:
        kwargs["response_format"] = {"type": "json_object"}
    for unsupported in list(cap.unsupported_parameters):
        kwargs.pop(unsupported, None)

    if nonstream_safe_mode:
        # OpenAI-compatible streaming was the crash source in macOS logs for a
        # plain "hi" prompt. Keep the same UI path by yielding the final text as
        # one chunk, and keep live tool loops opt-in while this path is validated.
        kwargs["stream"] = False
        kwargs.pop("tools", None)
        kwargs.pop("parallel_tool_calls", None)
        print("[llm] macOS OpenAI-compatible route using non-streaming safe mode", flush=True)
        text = _openai_compat_stdlib_completion_text(provider, kwargs)
        if text:
            if image_base64:
                _update_route_capabilities(provider, model, supports_images=True)
            yield text
        return

    def _retry_without_live_tools(reason: str) -> Generator[str, None, None]:
        """Handle retry without live tools for LLM clients client."""
        print(f"[llm] OpenAI-compatible live tools disabled for this route: {reason}", flush=True)
        _update_route_capabilities(provider, model, supports_tools=False)
        yield from _stream_openai_compat(
            user_message,
            image_base64,
            model,
            client,
            _inject_frontloaded_tool_context(
                ambient_context,
                allowed_tools,
                query=user_message,
            ),
            memory_context,
            use_tools=False,
            allowed_tools=allowed_tools,
            allow_screenshot_tool=False,
            screenshot_tool_b64=screenshot_tool_b64,
            provider=provider,
            max_tokens=max_tokens,
            temperature=temperature,
            json_mode=json_mode,
            history=history,
            system_prompt=system_prompt,
        )

    # Stream first round — yield text, accumulate any tool-call deltas.
    tool_calls_acc: dict[int, dict] = {}  # index → {id, name, arguments}
    finish_reason = None
    first_round_text: list[str] = []

    try:
        with client.chat.completions.create(**kwargs) as stream:
            _update_route_capabilities(provider, model, supports_stream=True)
            if image_base64:
                _update_route_capabilities(provider, model, supports_images=True)
            if tools:
                _update_route_capabilities(provider, model, supports_tools=True)
            for chunk in stream:
                if not chunk.choices:
                    continue
                choice = chunk.choices[0]
                finish_reason = choice.finish_reason or finish_reason
                delta = choice.delta
                if delta.content:
                    if tools:
                        first_round_text.append(delta.content)
                    else:
                        yield delta.content
                for tc in (delta.tool_calls or []):
                    entry = tool_calls_acc.setdefault(tc.index, {"id": "", "name": "", "arguments": ""})
                    if tc.id:
                        entry["id"] = tc.id
                    if tc.function:
                        if tc.function.name:
                            entry["name"] += tc.function.name
                        if tc.function.arguments:
                            entry["arguments"] += tc.function.arguments
    except Exception as exc:
        _record_route_error_capabilities(provider, model, exc)
        retry_kwargs = _without_unsupported_parameter(kwargs, exc)
        if retry_kwargs is not None:
            name = _unsupported_parameter_name(exc)
            _mark_unsupported_parameter(provider, model, name)
            if name == "tools":
                yield from _retry_without_live_tools(str(exc))
                return
            print(f"[llm] OpenAI-compatible rejected unsupported parameter {name!r}; retrying without it", flush=True)
            yield from _stream_openai_compat(
                user_message,
                image_base64,
                model,
                client,
                ambient_context,
                memory_context,
                use_tools=use_tools,
                allowed_tools=allowed_tools,
                pinned_tools=pinned_tools,
                allow_screenshot_tool=allow_screenshot_tool,
                screenshot_tool_b64=screenshot_tool_b64,
                provider=provider,
                max_tokens=max_tokens,
                temperature=temperature,
                json_mode=json_mode,
                history=history,
                system_prompt=system_prompt,
            )
            return
        if tools and _tools_not_supported_error(exc):
            yield from _retry_without_live_tools(str(exc))
            return
        if not kwargs.get("stream") or not _stream_mode_error(exc):
            raise
        if not _stream_false_required_error(exc):
            _update_route_capabilities(provider, model, supports_stream=False, requires_stream=False)
        retry_kwargs = dict(kwargs)
        retry_kwargs["stream"] = False
        print("[llm] OpenAI-compatible stream rejected; retrying non-streaming", flush=True)
        while True:
            try:
                response = client.chat.completions.create(**retry_kwargs)
                break
            except Exception as create_exc:
                _record_route_error_capabilities(provider, model, create_exc)
                next_kwargs = _without_unsupported_parameter(retry_kwargs, create_exc)
                if next_kwargs is not None:
                    _mark_unsupported_parameter(provider, model, _unsupported_parameter_name(create_exc))
                    retry_kwargs = next_kwargs
                    continue
                if tools and _tools_not_supported_error(create_exc):
                    yield from _retry_without_live_tools(str(create_exc))
                    return
                raise
        choice = response.choices[0] if getattr(response, "choices", None) else None
        finish_reason = getattr(choice, "finish_reason", None) or finish_reason
        text = _openai_compat_message_text(response)
        if text:
            if image_base64:
                _update_route_capabilities(provider, model, supports_images=True)
            if tools:
                first_round_text.append(text)
            else:
                yield text
        tool_calls_acc = _openai_compat_message_tool_calls(response)

    if finish_reason != "tool_calls" or not tool_calls_acc:
        first_round_joined = "".join(first_round_text)
        if tools and allow_screenshot_tool and _looks_like_screenshot_tool_request(first_round_joined):
            print("[llm] text requested capture_screen; continuing with implicit tool call", flush=True)
            tool_calls_acc = {
                0: {
                    "id": "implicit_capture_screen_1",
                    "name": "capture_screen",
                    "arguments": "{}",
                }
            }
        else:
            for text in first_round_text:
                yield text
            return

    # Build the assistant tool-call turn and append it to messages.
    assistant_tool_calls = [
        {"id": tc["id"], "type": "function",
         "function": {"name": tc["name"], "arguments": tc["arguments"]}}
        for tc in tool_calls_acc.values()
    ]
    messages.append({"role": "assistant", "content": None, "tool_calls": assistant_tool_calls})

    # Tool-call loop (up to 3 rounds, non-streaming for follow-ups).
    current_model = model
    vision_mode = False
    for _round in range(3):
        pending_image_b64: str | None = None
        for tc in tool_calls_acc.values():
            try:
                inputs = _json.loads(tc["arguments"] or "{}")
            except Exception:
                inputs = {}
            if tc["name"] == "capture_screen":
                # OpenAI tool messages can't carry an image, so acknowledge the
                # call with text and deliver the screenshot as a user message
                # below, then answer on a vision-capable model.
                print("[llm] tool_call capture_screen: resolving screenshot", flush=True)
                pending_image_b64 = _capture_screen_b64(screenshot_tool_b64)
                ack = (
                    "Screenshot captured; it is attached in the next message."
                    if pending_image_b64
                    else "Screen capture failed; no image is available."
                )
                messages.append({"role": "tool", "tool_call_id": tc["id"], "content": ack})
                _log_context("tool: capture_screen", "<screenshot>" if pending_image_b64 else "<failed>")
            else:
                result = _execute_model_tool(tc["name"], inputs, allowed_tools=allowed_tools)
                messages.append({"role": "tool", "tool_call_id": tc["id"], "content": result})

        if pending_image_b64:
            current_model = _openai_vision_model(provider, current_model)
            vision_mode = True  # answer with the image; drop tools for follow-ups
            print(f"[llm] tool_call capture_screen: follow-up vision model={current_model}", flush=True)
            messages.append({
                "role": "user",
                "content": [
                    {"type": "text", "text": "Here is the screenshot you requested."},
                    {"type": "image_url",
                     "image_url": {"url": f"data:image/png;base64,{pending_image_b64}"}},
                ],
            })

        follow_up_kwargs = {
            "model": current_model,
            "messages": messages,
            "stream": False,
        }
        _apply_sampling(follow_up_kwargs, current_model, 0.5 if temperature is None else temperature)
        if max_tokens != 0:
            _apply_max_output(follow_up_kwargs, current_model, max_tokens or 2048)
        if tools and not vision_mode:
            follow_up_kwargs["tools"] = tools
        follow_up = client.chat.completions.create(**follow_up_kwargs)
        choice = follow_up.choices[0]
        text = choice.message.content or ""
        if text:
            yield text
        if choice.finish_reason != "tool_calls":
            return
        # Another tool round — update accumulator and append assistant turn.
        tool_calls_acc = {}
        if choice.message.tool_calls:
            for i, tc in enumerate(choice.message.tool_calls):
                tool_calls_acc[i] = {
                    "id": tc.id,
                    "name": tc.function.name,
                    "arguments": tc.function.arguments or "",
                }
            messages.append({
                "role": "assistant", "content": None,
                "tool_calls": [
                    {"id": tc["id"], "type": "function",
                     "function": {"name": tc["name"], "arguments": tc["arguments"]}}
                    for tc in tool_calls_acc.values()
                ],
            })


# ------------------------------------------------------------------
# Codex (ChatGPT subscription) -” Responses API streaming
# ------------------------------------------------------------------

def _build_codex_text(user_message: str, ambient_context: str = "", memory_context: str = "") -> str:
    """Prepend context into a single text block for the Codex endpoint."""
    parts = []
    if memory_context:
        parts.append(memory_context)
    if ambient_context:
        parts.append(f"---\n{ambient_context}")
    parts.append(user_message)
    return "\n\n".join(parts)


def _codex_history_prefix(prior_turns: list[dict]) -> str:
    """Render sanitized history into the single-text Codex request format."""
    lines: list[str] = []
    for turn in prior_turns:
        role = str(turn.get("role") or "").strip().lower()
        content = str(turn.get("content") or "").strip()
        if not content:
            continue
        label = "User" if role == "user" else "Assistant"
        lines.append(f"{label}: {content}")
    return ("\n".join(lines) + "\n") if lines else ""


def _stream_codex(
    user_message: str,
    model: str,
    client,
    ambient_context: str = "",
    memory_context: str = "",
    use_tools: bool = False,
    allowed_tools: list[str] | None = None,
    pinned_tools: list[str] | None = None,
    history: list[dict] | None = None,
    system_prompt: str | None = None,
) -> Generator[str, None, None]:
    """Stream a response via the Codex endpoint using the Responses API."""
    prior_turns = _sanitize_history(history)
    history_prefix = _codex_history_prefix(prior_turns)
    text = _build_codex_text(
        (history_prefix + user_message) if history_prefix else user_message,
        ambient_context,
        memory_context,
    )
    if use_tools:
        tools = _get_responses_tool_schemas(
            user_message,
            include_general=True,
            allowed_tools=allowed_tools,
            pinned_tools=pinned_tools,
        )
        if tools:
            _log_offered_model_tools(
                "chatgpt",
                model,
                allowed_tools=allowed_tools,
                pinned_tools=pinned_tools,
                schemas=[{"function": {"name": t.get("name", "")}} for t in tools],
                openai_format=True,
            )
            try:
                yield from _run_responses_tool_loop(
                    client,
                    {
                        "model": model,
                        "input": [{"type": "message", "role": "user", "content": [{"type": "input_text", "text": text}]}],
                        "instructions": _with_tools_note(system_prompt or config.get_system_prompt(), True),
                        "tools": tools,
                        "store": False,
                    },
                    provider="chatgpt",
                    model=model,
                    allowed_tools=allowed_tools,
                )
                return
            except Exception as exc:
                if (
                    not _requires_stream_error(exc)
                    and not _tools_not_supported_error(exc)
                    and _without_unsupported_parameter({"tools": tools}, exc) is None
                ):
                    raise
                if _requires_stream_error(exc):
                    _update_route_capabilities("chatgpt", model, supports_stream=True, requires_stream=True)
                _update_route_capabilities("chatgpt", model, supports_tools=False)
                print(
                    f"[llm] ChatGPT/Codex live tools unavailable; front-loading supported context instead. {exc}",
                    flush=True,
                )
        else:
            print(
                f"[llm] ChatGPT/Codex tools requested, but no schemas were available. allowed={allowed_tools or []}",
                flush=True,
            )
        ambient_context = _inject_frontloaded_tool_context(
            ambient_context,
            allowed_tools,
            query=user_message,
        )
        text = _build_codex_text(
            (history_prefix + user_message) if history_prefix else user_message,
            ambient_context,
            memory_context,
        )
    yield from _response_stream_text(
        client,
        {
            "model": model,
            "input": [{"type": "message", "role": "user", "content": [{"type": "input_text", "text": text}]}],
            "instructions": system_prompt or config.get_system_prompt(),
            "store": False,
        },
        provider="chatgpt",
        model=model,
    )


def _stream_codex_vision(
    user_message: str,
    image_base64: str,
    model: str,
    client,
    ambient_context: str = "",
    memory_context: str = "",
) -> Generator[str, None, None]:
    """Stream a vision response via the Codex endpoint (Responses API with image input)."""
    text = _build_codex_text(user_message, ambient_context, memory_context)
    input_content = [
        {"type": "input_text",  "text": text},
        {"type": "input_image", "image_url": f"data:image/png;base64,{image_base64}"},
    ]
    for chunk in _response_stream_text(
        client,
        {
            "model": model,
            "input": [{"type": "message", "role": "user", "content": input_content}],
            "instructions": config.get_system_prompt(),
            "store": False,
        },
        provider="chatgpt",
        model=model,
    ):
        _update_route_capabilities("chatgpt", model, supports_images=True)
        yield chunk


# ------------------------------------------------------------------
# Anthropic Claude  -”  shared tool-loop helper
# ------------------------------------------------------------------

def _run_anthropic_tool_loop(
    client,
    messages: list,
    first_response,
    model: str,
    system: str,
    max_tokens: int,
    prompt: str = "",
    include_general: bool = True,
    include_screenshot: bool = False,
    screenshot_tool_b64: str | None = None,
    allowed_tools: list[str] | None = None,
    pinned_tools: list[str] | None = None,
) -> Generator[str, None, None]:
    """
    Execute Anthropic tool calls and yield text from subsequent rounds.
    *messages* is mutated in place (assistant + tool-result turns are appended).
    """
    messages.append({"role": "assistant", "content": first_response.content})
    final = first_response
    current_model = model
    for _round in range(3):
        tool_results = []
        for block in final.content:
            if getattr(block, "type", "") != "tool_use":
                continue
            if block.name == "capture_screen":
                # Return the screenshot as an image block; upgrade the loop to a
                # vision-capable model for the rounds that follow.
                print("[llm] tool_call capture_screen: resolving screenshot", flush=True)
                b64 = _capture_screen_b64(screenshot_tool_b64)
                if b64:
                    current_model = _anthropic_vision_model(current_model)
                    content = [{
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": "image/png",
                            "data": b64,
                        },
                    }]
                else:
                    content = "Screen capture failed; no image is available."
                _log_context("tool: capture_screen", "<screenshot>" if b64 else "<failed>")
            else:
                content = _execute_model_tool(block.name, block.input or {}, allowed_tools=allowed_tools)
            tool_results.append({
                "type": "tool_result",
                "tool_use_id": block.id,
                "content": content,
            })
        if not tool_results:
            return
        messages.append({"role": "user", "content": tool_results})
        print(f"[llm] tool_call follow-up: model={current_model}", flush=True)
        response = client.messages.create(
            model=current_model,
            max_tokens=max_tokens,
            system=system,
            messages=messages,
            tools=_get_tool_schemas(
                prompt,
                include_general=include_general,
                include_screenshot=include_screenshot,
                allowed_tools=allowed_tools,
                pinned_tools=pinned_tools,
            ),
        )
        for block in response.content:
            if getattr(block, "type", "") == "text":
                text = getattr(block, "text", "")
                if text:
                    yield text
        if response.stop_reason != "tool_use":
            return
        final = response
        messages.append({"role": "assistant", "content": response.content})


def _stream_anthropic(
    user_message: str,
    image_base64: str | None,
    model: str,
    client,
    ambient_context: str = "",
    memory_context: str = "",
    use_tools: bool = False,
    allowed_tools: list[str] | None = None,
    pinned_tools: list[str] | None = None,
    allow_screenshot_tool: bool = False,
    screenshot_tool_b64: str | None = None,
    max_tokens: int | None = None,
    temperature: float | None = None,
    history: list[dict] | None = None,
    system_prompt: str | None = None,
) -> Generator[str, None, None]:
    """Stream anthropic."""
    tools_active = use_tools or allow_screenshot_tool
    prior_turns = _sanitize_history(history)
    cap = _get_route_capabilities("anthropic", model)
    if tools_active and cap.supports_tools is False:
        print(
            f"[llm] Anthropic live tools disabled by route capability cache; "
            f"front-loading supported context instead. allowed={allowed_tools or []}",
            flush=True,
        )
        ambient_context = _inject_frontloaded_tool_context(
            ambient_context,
            allowed_tools,
            query=user_message,
        )
        tools_active = False
        use_tools = False
        allow_screenshot_tool = False

    system = _with_screenshot_note(system_prompt or config.get_system_prompt(), allow_screenshot_tool)
    system = _with_tools_note(system, use_tools)
    system = _with_memory_search_note(system, allowed_tools if use_tools else None)
    system = _with_memory_save_note(system, allowed_tools if use_tools else None)
    if memory_context:
        system += f"\n\n{memory_context}"
    if ambient_context:
        system += f"\n\n---\n{ambient_context}"

    # Anthropic requires an explicit max_tokens, so "no app cap" (0) maps to a
    # generous per-request ceiling rather than being truly unlimited.
    anthropic_max_tokens = _ANTHROPIC_UNCAPPED_MAX_TOKENS if max_tokens == 0 else (max_tokens or 2048)
    anthropic_tool_max_tokens = _ANTHROPIC_UNCAPPED_MAX_TOKENS if max_tokens == 0 else (max_tokens or 512)

    if image_base64:
        content = [
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/png",
                    "data": image_base64,
                },
            },
            {"type": "text", "text": user_message},
        ]
    else:
        content = user_message

    # --- No tools: original streaming path (lowest latency) ---
    if not tools_active:
        request = {
            "model": model,
            "max_tokens": anthropic_max_tokens,
            "system": system,
            "messages": [*prior_turns, {"role": "user", "content": content}],
        }
        if temperature is not None:
            request["temperature"] = temperature
        with client.messages.stream(**request) as stream:
            _update_route_capabilities("anthropic", model, supports_stream=True)
            if image_base64:
                _update_route_capabilities("anthropic", model, supports_images=True)
            for text in stream.text_stream:
                yield text
        return

    # --- Tool-enabled path: stream first round for fast first-token ---
    # If no tool is called (common case), text streams immediately.
    # Only falls back to blocking create() if Claude actually invokes a tool.
    messages: list[dict] = [*prior_turns, {"role": "user", "content": content}]
    tool_schemas = _get_tool_schemas(
        user_message,
        include_general=use_tools,
        include_screenshot=allow_screenshot_tool,
        allowed_tools=allowed_tools,
        pinned_tools=pinned_tools,
    )
    _log_offered_model_tools(
        "anthropic",
        model,
        allowed_tools=allowed_tools,
        pinned_tools=pinned_tools,
        schemas=tool_schemas,
    )
    if not tool_schemas:
        yield from _stream_anthropic(
            user_message,
            image_base64,
            model,
            client,
            _inject_frontloaded_tool_context(
                ambient_context,
                allowed_tools,
                query=user_message,
            ),
            memory_context,
            use_tools=False,
            allowed_tools=allowed_tools,
            allow_screenshot_tool=False,
            screenshot_tool_b64=screenshot_tool_b64,
            max_tokens=max_tokens,
            temperature=temperature,
            history=history,
            system_prompt=system_prompt,
        )
        return

    request = {
        "model": model,
        "max_tokens": anthropic_max_tokens,
        "system": system,
        "messages": messages,
        "tools": tool_schemas,
    }
    if temperature is not None:
        request["temperature"] = temperature
    first_round_text: list[str] = []
    try:
        with client.messages.stream(**request) as stream:
            _update_route_capabilities("anthropic", model, supports_stream=True, supports_tools=True)
            if image_base64:
                _update_route_capabilities("anthropic", model, supports_images=True)
            for text in stream.text_stream:
                first_round_text.append(text)
            final = stream.get_final_message()
    except Exception as exc:
        _record_route_error_capabilities("anthropic", model, exc)
        if _tools_not_supported_error(exc):
            print("[llm] Anthropic live tools rejected; retrying with front-loaded context", flush=True)
            _update_route_capabilities("anthropic", model, supports_tools=False)
            yield from _stream_anthropic(
                user_message,
                image_base64,
                model,
                client,
                _inject_frontloaded_tool_context(
                    ambient_context,
                    allowed_tools,
                    query=user_message,
                ),
                memory_context,
                use_tools=False,
                allowed_tools=allowed_tools,
                allow_screenshot_tool=False,
                screenshot_tool_b64=screenshot_tool_b64,
                max_tokens=max_tokens,
                temperature=temperature,
                history=history,
                system_prompt=system_prompt,
            )
            return
        if _streaming_not_supported_error(exc):
            _update_route_capabilities("anthropic", model, supports_stream=False, requires_stream=False)
        raise

    if final.stop_reason != "tool_use":
        for text in first_round_text:
            yield text
        return

    # A tool was called -” execute it and do followup round(s) non-streaming.
    yield from _run_anthropic_tool_loop(
        client,
        messages,
        final,
        model,
        system,
        max_tokens=anthropic_tool_max_tokens,
        prompt=user_message,
        include_general=use_tools,
        include_screenshot=allow_screenshot_tool,
        screenshot_tool_b64=screenshot_tool_b64,
        allowed_tools=allowed_tools,
        pinned_tools=pinned_tools,
    )


# ------------------------------------------------------------------
# Inline rewrite / fix  (Ctrl+Shift+Q)
# ------------------------------------------------------------------

def stream_rewrite(selected_text: str, intent_prompt: str = "Rewrite or fix the following text") -> Generator[str, None, None]:
    """
    Stream a rewrite/fix of the selected text using the primary LLM.

    The system prompt instructs the model to output raw replacement text only -”
    no prose, no markdown, no explanation.  The result is pasted back directly.

    Args:
        selected_text:  The text to rewrite.
        intent_prompt:  The instruction (e.g. "Fix the grammar and spelling").
                        Taken from the caller's chosen intent row.
    """
    import time
    ts = time.strftime("%H:%M:%S")
    print(f"[llm {ts}] Rewrite request ({len(selected_text)} chars) -” {intent_prompt[:60]!r}")
    user_message = f"{intent_prompt}:\n\n{selected_text}"
    candidates = _route_candidates(config.LLM_PROVIDER, config.LLM_MODEL, config.LLM_FALLBACKS)
    yield from _stream_with_fallbacks(
        "rewrite",
        candidates,
        lambda provider, model: _stream_single_rewrite_route(provider, model, user_message),
    )


def _stream_single_rewrite_route(provider: str, model: str, user_message: str) -> Generator[str, None, None]:
    """Stream single rewrite route."""
    _check_route_config(provider, model, "LLM")
    _log_model_route("rewrite", provider, model, use_tools=False)
    if provider in _OPENAI_COMPAT_PROVIDER_SET:
        kwargs = {
            "model": model,
            "messages": [
                {"role": "system", "content": _REWRITE_SYSTEM_PROMPT},
                {"role": "user",   "content": user_message},
            ],
            "stream": not _use_macos_openai_compat_non_streaming(provider),
        }
        _apply_sampling(kwargs, model, 0.3)
        _apply_max_output(kwargs, model, 1024)
        yield from _stream_openai_compat_plain(provider, model, kwargs)
        return
    elif provider == "anthropic":
        client = _dynamic_anthropic_client()
        request = {
            "model": model,
            "max_tokens": 1024,
            "system": _REWRITE_SYSTEM_PROMPT,
            "messages": [{"role": "user", "content": user_message}],
        }
        _apply_sampling(request, model, 0.3)
        with client.messages.stream(**request) as stream:
            for text in stream.text_stream:
                yield text
    elif provider == "copilot":
        from core.auth import copilot_client
        yield from copilot_client.stream(
            user_message,
            model,
            system=_REWRITE_SYSTEM_PROMPT,
        )
    elif provider == "chatgpt":
        yield from _response_stream_text(
            _get_codex_client(),
            {
                "model": model,
                "input": [{"type": "message", "role": "user", "content": [{"type": "input_text", "text": user_message}]}],
                "instructions": _REWRITE_SYSTEM_PROMPT,
                "store": False,
            },
            provider="chatgpt",
            model=model,
        )
    else:
        raise ValueError(f"Unknown rewrite provider: {provider}")


# ------------------------------------------------------------------
# Multi-turn (chat window)
# ------------------------------------------------------------------

def _history_user_has_image(messages: list) -> bool:
    """True if any user turn carries an attached screenshot."""
    return any(
        m.get("role") == "user" and m.get("image_base64")
        for m in messages
    )


def _history_text_payload(messages: list) -> tuple[str, str, list[dict]]:
    """Return (system, latest user message, prior text turns) for shared chat routing."""
    system = next(
        (
            str(m.get("content") or "")
            for m in messages
            if m.get("role") == "system"
        ),
        config.get_system_prompt(),
    )
    turns: list[dict[str, str]] = []
    for m in messages:
        role = str(m.get("role") or "").strip().lower()
        if role not in {"user", "assistant"}:
            continue
        content = m.get("content")
        if not isinstance(content, str):
            continue
        turns.append({"role": role, "content": content})
    last_user_idx = next(
        (idx for idx in range(len(turns) - 1, -1, -1) if turns[idx].get("role") == "user"),
        -1,
    )
    if last_user_idx < 0:
        return system, "", turns
    return system, turns[last_user_idx].get("content") or "", turns[:last_user_idx]


def _openai_history_payload(messages: list) -> list:
    """Render history for the OpenAI-compatible chat completions API.

    User turns that carry a screenshot become a multimodal content array so the
    image is replayed on every turn; everything else stays a plain string."""
    out: list[dict] = []
    for m in messages:
        role = m.get("role")
        text = str(m.get("content") or "")
        image = m.get("image_base64") if role == "user" else None
        if image:
            content: object = [
                {"type": "text", "text": text},
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{image}"},
                },
            ]
        else:
            content = text
        out.append({"role": role, "content": content})
    return out


def _anthropic_history_payload(messages: list) -> tuple[str, list]:
    """Render history as (system, turns) for the Anthropic messages API.

    User turns with a screenshot become image+text content blocks."""
    system = next((m["content"] for m in messages if m.get("role") == "system"), "")
    turns: list[dict] = []
    for m in messages:
        role = m.get("role")
        if role == "system":
            continue
        text = str(m.get("content") or "")
        image = m.get("image_base64") if role == "user" else None
        if image:
            content: object = [
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": "image/png",
                        "data": image,
                    },
                },
                {"type": "text", "text": text},
            ]
        else:
            content = text
        turns.append({"role": role, "content": content})
    return system, turns


_CACHE_CONTROL = {"type": "ephemeral"}


def _anthropic_cached_system(system: str):
    """Render the system prompt as a cacheable block.

    Anthropic does not cache anything unless you opt in with a ``cache_control``
    breakpoint. The system prompt is byte-identical every turn, so caching it
    (together with the tool list, which renders just before it) is a safe win.
    """
    text = system or ""
    if not text:
        return text
    return [{"type": "text", "text": text, "cache_control": dict(_CACHE_CONTROL)}]


def _mark_anthropic_history_cache(turns: list) -> None:
    """Put a cache breakpoint at the end of the stable history, in place.

    ``turns[-1]`` is the volatile latest user turn (new question + per-turn
    memory); everything before it is frozen conversation history. Marking the
    block just before the latest turn lets the next request reuse the whole prior
    prefix — including any replayed screenshot — at cache-read pricing.
    """
    if len(turns) < 2:
        return  # only the latest turn exists; nothing stable to cache yet
    target = turns[-2]
    content = target.get("content")
    if isinstance(content, str):
        target["content"] = [
            {"type": "text", "text": content, "cache_control": dict(_CACHE_CONTROL)}
        ]
    elif isinstance(content, list) and content:
        last = dict(content[-1])
        last["cache_control"] = dict(_CACHE_CONTROL)
        content[-1] = last


def stream_response_with_history(
    messages: list,
    memory_context: str = "",
    use_tools: bool = False,
    allowed_tools: list[str] | None = None,
    pinned_tools: list[str] | None = None,
) -> Generator[str, None, None]:
    """
    Stream a response given a pre-built messages list including history.
    Uses CHAT_LLM_PROVIDER / CHAT_LLM_MODEL (defaults to LLM_PROVIDER / LLM_MODEL).

    When any user turn carries an ``image_base64`` screenshot, the whole
    conversation is routed through VISION_LLM_PROVIDER / MODEL (like the one-shot
    vision path) and the image is replayed on every turn, so the model keeps
    seeing it on follow-up questions.

    Args:
        messages:        [{{"role": "system"|"user"|"assistant", "content": str,
                         "image_base64": str?}}, ...]
        memory_context:  Pre-formatted LTM facts from core.memory -” attached to
                         the latest user turn so the model is aware of user facts.
    """
    # Attach memory to the latest user turn, NOT the system prompt. Memory is
    # re-retrieved every turn, so folding it into the frozen system prompt would
    # change the very front of the prompt each time and invalidate prompt caching
    # for the whole conversation (system + history render before it). Keeping it
    # on the volatile last user turn lets the stable prefix stay byte-identical so
    # providers can cache the replayed history (including any screenshot).
    if memory_context:
        messages = [dict(m) for m in messages]  # copy -” don't mutate the caller's turns
        for i in range(len(messages) - 1, -1, -1):
            if messages[i].get("role") == "user":
                base = messages[i].get("content")
                if isinstance(base, str):
                    messages[i]["content"] = f"{base}\n\n{memory_context}".strip()
                break
        else:
            # No user turn — fall back to a leading system message.
            messages = [{"role": "system", "content": memory_context}] + messages
    has_image = _history_user_has_image(messages)
    if has_image:
        kind, route_name = "chat-vision", "VISION_LLM"
        candidates = _route_candidates(
            config.VISION_LLM_PROVIDER,
            config.VISION_LLM_MODEL,
            config.VISION_LLM_FALLBACKS,
        )
    else:
        kind, route_name = "chat", "CHAT_LLM"
        candidates = _route_candidates(
            config.CHAT_LLM_PROVIDER,
            config.CHAT_LLM_MODEL,
            config.CHAT_LLM_FALLBACKS,
        )
    yield from _stream_with_fallbacks(
        kind,
        candidates,
        lambda provider, model: _stream_single_history_route(
            provider,
            model,
            messages,
            route_name=route_name,
            use_tools=use_tools,
            allowed_tools=allowed_tools,
            pinned_tools=pinned_tools,
        ),
    )


def _stream_single_history_route(
    provider: str,
    model: str,
    messages: list,
    route_name: str = "CHAT_LLM",
    use_tools: bool = False,
    allowed_tools: list[str] | None = None,
    pinned_tools: list[str] | None = None,
) -> Generator[str, None, None]:
    """Stream single history route."""
    if not _history_user_has_image(messages):
        system_msg, user_message, history = _history_text_payload(messages)
        chat_temperature = 0.7 if provider in _OPENAI_COMPAT_PROVIDER_SET else None
        yield from _stream_single_response_route(
            provider,
            model,
            user_message,
            None,
            "",
            "",
            use_tools=use_tools,
            route_name=route_name,
            allowed_tools=allowed_tools,
            pinned_tools=pinned_tools,
            max_tokens=1024,
            temperature=chat_temperature,
            history=history,
            system_prompt=system_msg,
            route_kind="chat",
        )
        return

    _check_route_config(provider, model, route_name)
    _log_model_route("chat", provider, model, use_tools=use_tools)
    if provider in _OPENAI_COMPAT_PROVIDER_SET:
        kwargs = {
            "model": model,
            "messages": _openai_history_payload(messages),
            "stream": not _use_macos_openai_compat_non_streaming(provider),
        }
        _apply_sampling(kwargs, model, 0.7)
        _apply_max_output(kwargs, model, 1024)
        yield from _stream_openai_compat_plain(provider, model, kwargs)
        return
    elif provider == "anthropic":
        client = _dynamic_anthropic_client()
        # Build turns with screenshots replayed as image blocks; drops any keys
        # the API would reject (e.g. image_base64) by reconstructing each turn.
        system, turns = _anthropic_history_payload(messages)
        # Send the full, stable tool set (not the per-turn keyword-filtered one):
        # tools render at the front of the prompt, so a set that changes between
        # turns would invalidate the whole prompt cache. The fixed set keeps the
        # prefix byte-identical and stays cached with the system prompt.
        chat_tools = (
            _get_tool_schemas(
                allowed_tools=allowed_tools,
                pinned_tools=pinned_tools,
                unfiltered=True,
            )
            if use_tools
            else []
        )
        # Prompt caching: cache the frozen system prompt and the stable history
        # prefix so replayed turns (incl. screenshots) bill at ~0.1x on follow-ups.
        system = _anthropic_cached_system(system)
        _mark_anthropic_history_cache(turns)
        with client.messages.stream(
            model=model,
            max_tokens=1024,
            system=system,
            messages=turns,
            **({"tools": chat_tools} if chat_tools else {}),
        ) as stream:
            for text in stream.text_stream:
                yield text
            final = stream.get_final_message()

        if final.stop_reason != "tool_use":
            return

        yield from _run_anthropic_tool_loop(client, turns, final, model, system, max_tokens=1024)
    elif provider == "chatgpt":
        # For the chat history path, extract the last user message and use Responses API.
        # The full history is packed into a single user turn with prior turns prefixed.
        system_msg = next((m["content"] for m in messages if m["role"] == "system"), "")
        turns = [m for m in messages if m["role"] != "system"]
        # Build a text representation of prior context for the model
        last_user = turns[-1]["content"] if turns else ""
        history_prefix = ""
        for m in turns[:-1]:
            label = "User" if m["role"] == "user" else "Assistant"
            history_prefix += f"{label}: {m['content']}\n"
        full_input = (history_prefix + last_user) if history_prefix else last_user
        # Replay any screenshots in the conversation as input_image blocks so the
        # Responses model can still see them on follow-up turns.
        content: list[dict] = [{"type": "input_text", "text": full_input}]
        for m in turns:
            if m.get("role") == "user" and m.get("image_base64"):
                content.append({
                    "type": "input_image",
                    "image_url": f"data:image/png;base64,{m['image_base64']}",
                })
        if use_tools:
            tools = _get_responses_tool_schemas(
                str(last_user or full_input),
                include_general=True,
                allowed_tools=allowed_tools,
                pinned_tools=pinned_tools,
            )
            if tools:
                _log_offered_model_tools(
                    "chatgpt",
                    model,
                    allowed_tools=allowed_tools,
                    pinned_tools=pinned_tools,
                    schemas=[{"function": {"name": t.get("name", "")}} for t in tools],
                    openai_format=True,
                )
                try:
                    instructions = _with_tools_note(system_msg, True)
                    instructions = _with_memory_search_note(instructions, allowed_tools)
                    instructions = _with_memory_save_note(instructions, allowed_tools)
                    yield from _run_responses_tool_loop(
                        _get_chat_codex_client(),
                        {
                            "model": model,
                            "input": [{"type": "message", "role": "user", "content": content}],
                            "instructions": instructions,
                            "tools": tools,
                            "store": False,
                        },
                        provider="chatgpt",
                        model=model,
                        allowed_tools=allowed_tools,
                    )
                    return
                except Exception as exc:
                    if (
                        not _requires_stream_error(exc)
                        and not _tools_not_supported_error(exc)
                        and _without_unsupported_parameter({"tools": tools}, exc) is None
                    ):
                        raise
                    if _requires_stream_error(exc):
                        _update_route_capabilities("chatgpt", model, supports_stream=True, requires_stream=True)
                    _update_route_capabilities("chatgpt", model, supports_tools=False)
                    print(
                        f"[llm] ChatGPT/Codex chat tools unavailable; retrying without live tools. {exc}",
                        flush=True,
                    )
        yield from _response_stream_text(
            _get_chat_codex_client(),
            {
                "model": model,
                "input": [{"type": "message", "role": "user", "content": content}],
                "instructions": system_msg,
                "store": False,
            },
            provider="chatgpt",
            model=model,
        )
    elif provider == "copilot":
        from core.auth import copilot_client

        system_msg = next((m["content"] for m in messages if m["role"] == "system"), "")
        turns = [m for m in messages if m["role"] != "system"]
        full_input = ""
        for m in turns:
            label = "User" if m["role"] == "user" else "Assistant"
            full_input += f"{label}: {m['content']}\n"
        yield from copilot_client.stream(
            full_input.strip(),
            model,
            system=system_msg,
            session_id="wisp-chat",
        )
    else:
        raise ValueError(f"Unknown chat LLM provider: {provider}")
